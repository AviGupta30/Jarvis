"""
ppt_tool.py — Premium Visual-First PPT Engine v5
=================================================
Design philosophy:
  • Visuals are DOMINANT — image gets 58%+ of the slide width/height
  • Background has depth — gradient strips, decorative geometry
  • Cards are content-dense — colored numbered tags, bold sub-headers, tight spacing
  • Image frame is premium — thick border + L-corner accents + subtle inner glow ring
  • Every slide has a polished header: brand chip + centered title box + accent line
  • Bottom: full-width bar + slide counter
"""
from __future__ import annotations
import math
import io, json, re, time
from dataclasses import dataclass
from pathlib import Path
from typing import Optional, Generator

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    import subprocess, sys
    subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx", "-q"], check=False)
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

try:
    from groq import Groq
    from app.core.config import settings
    _GROQ = Groq(api_key=settings.GROQ_API_KEY)
except Exception:
    _GROQ = None

_LAST_DECK_PATH: Optional[str] = None
_LAST_PLAN:      Optional[dict] = None

W = Inches(13.33)
H = Inches(7.5)

# ─── Premium Palettes ──────────────────────────────────────────────────────────
PERSONALITIES = {
    "ocean_pro": {
        "name": "Ocean Pro", "desc": "Deep navy + electric cyan.",
        "bg": "050E1F", "bg2": "09162E", "card": "0D1F3C", "card2": "0A1830",
        "text": "D8EEFF", "sub": "7AACCC", "ac1": "00C8FF", "ac2": "0077CC", "ac3": "00FFB3",
        "border": "00C8FF", "hdr_bg": "00C8FF", "hdr_text": "050E1F", "bar": "0077CC",
        "tag_colors": ["00C8FF", "0077CC", "00FFB3", "4488FF"],
    },
    "neon_dark": {
        "name": "Neon Dark", "desc": "Black + vivid purple/pink.",
        "bg": "06020D", "bg2": "0E0520", "card": "140830", "card2": "1A0A3A",
        "text": "F5EEFF", "sub": "C4A8EE", "ac1": "BF5FFF", "ac2": "FF2D78", "ac3": "00F5FF",
        "border": "BF5FFF", "hdr_bg": "BF5FFF", "hdr_text": "06020D", "bar": "FF2D78",
        "tag_colors": ["BF5FFF", "FF2D78", "00F5FF", "FF8800"],
    },
    "emerald": {
        "name": "Emerald", "desc": "Dark forest green + gold.",
        "bg": "021A0C", "bg2": "042A14", "card": "083820", "card2": "0A4228",
        "text": "D8FFF0", "sub": "88CCAA", "ac1": "00E878", "ac2": "00AA55", "ac3": "FFD700",
        "border": "00E878", "hdr_bg": "00E878", "hdr_text": "021A0C", "bar": "00AA55",
        "tag_colors": ["00E878", "00AA55", "FFD700", "44DDBB"],
    },
    "clean_light": {
        "name": "Clean Light", "desc": "White + corporate blue.",
        "bg": "F0F5FF", "bg2": "E4EDFF", "card": "FFFFFF", "card2": "EBF2FF",
        "text": "0A1628", "sub": "3A5888", "ac1": "1A56DB", "ac2": "1040BB", "ac3": "F59E0B",
        "border": "1A56DB", "hdr_bg": "1A56DB", "hdr_text": "FFFFFF", "bar": "1040BB",
        "tag_colors": ["1A56DB", "F59E0B", "10B981", "EF4444"],
    },
    "synthwave": {
        "name": "Synthwave", "desc": "Neon pink + purple. Retro cyber.",
        "bg": "0A0118", "bg2": "130228", "card": "1C0535", "card2": "240645",
        "text": "FFE8FF", "sub": "D4A0D8", "ac1": "FF1493", "ac2": "9B30FF", "ac3": "00FFEE",
        "border": "FF1493", "hdr_bg": "FF1493", "hdr_text": "FFFFFF", "bar": "9B30FF",
        "tag_colors": ["FF1493", "9B30FF", "00FFEE", "FFAA00"],
    },
    "aurora": {
        "name": "Aurora", "desc": "Dark teal + bright green.",
        "bg": "011215", "bg2": "021C20", "card": "042C30", "card2": "05383D",
        "text": "D8FFF8", "sub": "80C8C0", "ac1": "00E8CC", "ac2": "00A898", "ac3": "80FF44",
        "border": "00E8CC", "hdr_bg": "00E8CC", "hdr_text": "011215", "bar": "00A898",
        "tag_colors": ["00E8CC", "00A898", "80FF44", "FFCC00"],
    },
    "hacker_terminal": {
        "name": "Hacker Terminal", "desc": "Pitch black + lime green.",
        "bg": "000000", "bg2": "040804", "card": "080D08", "card2": "0D140D",
        "text": "B8FFB8", "sub": "5EA85E", "ac1": "39FF14", "ac2": "20C000", "ac3": "FFFF00",
        "border": "39FF14", "hdr_bg": "39FF14", "hdr_text": "000000", "bar": "20C000",
        "tag_colors": ["39FF14", "20C000", "FFFF00", "FF8800"],
    },
    "solar_flare": {
        "name": "Solar Flare", "desc": "Charcoal + amber. Startup energy.",
        "bg": "0A0700", "bg2": "160E00", "card": "241500", "card2": "2E1A00",
        "text": "FFF8E8", "sub": "D4B870", "ac1": "FF8C00", "ac2": "FF5500", "ac3": "FFD700",
        "border": "FF8C00", "hdr_bg": "FF8C00", "hdr_text": "0A0700", "bar": "FF5500",
        "tag_colors": ["FF8C00", "FF5500", "FFD700", "FF2200"],
    },
    "midnight_exec": {
        "name": "Midnight Executive", "desc": "Deep indigo + platinum.",
        "bg": "05061A", "bg2": "090C28", "card": "101538", "card2": "161C44",
        "text": "E8EAF8", "sub": "8890C8", "ac1": "8080FF", "ac2": "5050DD", "ac3": "C0C0FF",
        "border": "8080FF", "hdr_bg": "8080FF", "hdr_text": "05061A", "bar": "5050DD",
        "tag_colors": ["8080FF", "5050DD", "C0C0FF", "FF8080"],
    },
    "cyber_dark": {
        "name": "Cyber Dark", "desc": "Dark navy + electric teal.",
        "bg": "030B18", "bg2": "060F25", "card": "091830", "card2": "0C1E3A",
        "text": "C8F0FF", "sub": "5AAABB", "ac1": "00F3FF", "ac2": "0088FF", "ac3": "00FF88",
        "border": "00F3FF", "hdr_bg": "00F3FF", "hdr_text": "030B18", "bar": "0088FF",
        "tag_colors": ["00F3FF", "0088FF", "00FF88", "FF8800"],
    },
    "arctic_clean": {
        "name": "Arctic Clean", "desc": "Ice white + deep teal.",
        "bg": "EDF8FA", "bg2": "D8EFF5", "card": "FFFFFF", "card2": "D0EBF5",
        "text": "001828", "sub": "005070", "ac1": "0088AA", "ac2": "006688", "ac3": "00CCAA",
        "border": "0088AA", "hdr_bg": "0088AA", "hdr_text": "FFFFFF", "bar": "006688",
        "tag_colors": ["0088AA", "006688", "00CCAA", "FF8800"],
    },
    "forest_calm": {
        "name": "Forest Calm", "desc": "Forest green + cream.",
        "bg": "041008", "bg2": "071A0C", "card": "0C2A12", "card2": "103416",
        "text": "D8F0DC", "sub": "80B888", "ac1": "44CC66", "ac2": "2EA050", "ac3": "CCFF44",
        "border": "44CC66", "hdr_bg": "44CC66", "hdr_text": "041008", "bar": "2EA050",
        "tag_colors": ["44CC66", "2EA050", "CCFF44", "FF8800"],
    },
    "ocean_gradient": {
        "name": "Ocean Gradient", "desc": "Deep ocean + cerulean.",
        "bg": "020B18", "bg2": "04102A", "card": "081830", "card2": "0B1E3A",
        "text": "C8E8FF", "sub": "6899BB", "ac1": "0099DD", "ac2": "0066BB", "ac3": "00DDCC",
        "border": "0099DD", "hdr_bg": "0099DD", "hdr_text": "020B18", "bar": "0066BB",
        "tag_colors": ["0099DD", "0066BB", "00DDCC", "FF8800"],
    },
    "velvet_noir": {
        "name": "Velvet Noir", "desc": "Deep plum + rose gold.",
        "bg": "0E0618", "bg2": "160828", "card": "1E0C3A", "card2": "280E48",
        "text": "F8E8FF", "sub": "C898D8", "ac1": "CC88FF", "ac2": "AA44CC", "ac3": "FFAA80",
        "border": "CC88FF", "hdr_bg": "CC88FF", "hdr_text": "0E0618", "bar": "AA44CC",
        "tag_colors": ["CC88FF", "AA44CC", "FFAA80", "FF2D78"],
    },
    "charcoal_minimal": {
        "name": "Charcoal Minimal", "desc": "True black + gold.",
        "bg": "0A0A0A", "bg2": "121212", "card": "1C1C1C", "card2": "242424",
        "text": "F0F0F0", "sub": "909090", "ac1": "CCA020", "ac2": "EEC040", "ac3": "FFFFFF",
        "border": "CCA020", "hdr_bg": "CCA020", "hdr_text": "0A0A0A", "bar": "EEC040",
        "tag_colors": ["CCA020", "EEC040", "FFFFFF", "FF5500"],
    },
}

# ─── Style Profiles ────────────────────────────────────────────────────────────
# Controls the entire visual DNA, not just colors.
STYLE_PROFILES = {
    "hackathon": {
        "name": "Hackathon / Tech Pitch",
        "desc": "Bold, dark, visual-dominant. Numbered badges, corner brackets, gradient strips.",
        # Decoration flags
        "corner_brackets": True,
        "numbered_badges": True,
        "tech_dots": True,
        "gradient_strips": True,
        "right_edge_panel": True,
        "drop_shadows": True,
        "decorative_circles": True,
        "colored_borders": True,
        "heavy_title_card": True,
        "filled_header": True,
        # Layout geometry
        "image_ratio": 0.57,
        "text_ratio": 0.41,
        "header_style": "pill",
        "bullet_style": "card",
        "card_border_width": 1.5,
        # Typography
        "title_font_size": 48,
        "header_font_size": 22,
        "body_font_size": 11,
        "title_font": "Calibri",
        "body_font": "Calibri",
        # Palette constraints
        "default_palette": "ocean_pro",
        "allowed_palettes": [
            "ocean_pro", "neon_dark", "emerald", "synthwave", "aurora",
            "hacker_terminal", "solar_flare", "cyber_dark", "velvet_noir",
        ],
    },
    "general": {
        "name": "General Purpose / Professional",
        "desc": "Clean, minimal, professional. No tech decorations, classic formatting.",
        # Decoration flags — ALL OFF
        "corner_brackets": False,
        "numbered_badges": False,
        "tech_dots": False,
        "gradient_strips": False,
        "right_edge_panel": False,
        "drop_shadows": False,
        "decorative_circles": False,
        "colored_borders": False,
        "heavy_title_card": False,
        "filled_header": False,
        # Layout geometry
        "image_ratio": 0.48,
        "text_ratio": 0.50,
        "header_style": "bar",
        "bullet_style": "inline",
        "card_border_width": 0.75,
        # Typography
        "title_font_size": 44,
        "header_font_size": 20,
        "body_font_size": 12,
        "title_font": "Calibri",
        "body_font": "Calibri",
        # Palette constraints
        "default_palette": "clean_light",
        "allowed_palettes": [
            "clean_light", "arctic_clean", "charcoal_minimal",
            "midnight_exec", "forest_calm", "ocean_gradient",
        ],
    },
}

_HACKATHON_KEYWORDS = [
    "hackathon", "pitch", "startup", "demo", "prototype",
    "mvp", "product launch", "tech demo", "innovation challenge",
    "pitch deck", "investor", "shark tank",
]

def _detect_purpose(prompt: str) -> str:
    """Auto-detect presentation purpose from the user prompt."""
    lower = prompt.lower()
    if any(kw in lower for kw in _HACKATHON_KEYWORDS):
        return "hackathon"
    return "general"


def _c(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:], 16))

def _bg_fill(slide, col: str):
    b = slide.background; b.fill.solid(); b.fill.fore_color.rgb = _c(col)

def _shape(slide, shp_type, l, t, w, h, fill=None, line=None, lw=1.0):
    sh = slide.shapes.add_shape(shp_type, l, t, w, h)
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = _c(fill)
    else: sh.fill.background()
    if line: sh.line.color.rgb = _c(line); sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    return sh

def _rect(slide, l, t, w, h, fill=None, line=None, lw=1.0):
    return _shape(slide, MSO_SHAPE.RECTANGLE, l, t, w, h, fill=fill, line=line, lw=lw)

def _round(slide, l, t, w, h, fill=None, line=None, lw=1.0):
    return _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, fill=fill, line=line, lw=lw)

def _oval(slide, l, t, w, h, fill=None, line=None, lw=1.0):
    return _shape(slide, MSO_SHAPE.OVAL, l, t, w, h, fill=fill, line=line, lw=lw)

def _tb(slide, text: str, l, t, w, h, font="Calibri", sz=11, bold=False, italic=False,
        col="FFFFFF", align=PP_ALIGN.LEFT, v_align=None):
    if not text: return
    
    from pptx.enum.shapes import MSO_SHAPE
    try:
        from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
    except ImportError:
        MSO_AUTO_SIZE = None
        MSO_ANCHOR = None
        
    tb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    tb.fill.background()
    tb.line.fill.background()
    
    tf = tb.text_frame
    tf.word_wrap = True
    if MSO_AUTO_SIZE:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if MSO_ANCHOR:
        tf.vertical_anchor = v_align if v_align else MSO_ANCHOR.TOP
        
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.name = font; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = _c(col)
    return tb

def _clean_image_path(image_path: str) -> str:
    if not image_path: return ""
    s = str(image_path)
    m = re.search(r'\[ATTACHED_FILE:\s*(.+?)(?:\s*\|.*)?\]', s)
    if m: return m.group(1).strip()
    return s.split('|')[0].strip()

def _get_image_aspect_ratio(path: str) -> float:
    """Returns width / height of image using PIL if available, else 1.78 (16:9)"""
    try:
        from PIL import Image
        with Image.open(path) as img:
            return img.width / max(img.height, 1)
    except Exception:
        return 1.78


# ─── Fluid Slot Geometry ───────────────────────────────────────────────────────
@dataclass
class SlotGeometry:
    """Computed image + text zone dimensions (in EMU — python-pptx native)."""
    img_x: float; img_y: float; img_w: float; img_h: float
    txt_x: float; txt_y: float; txt_w: float; txt_h: float


def compute_split_geometry(top: float, avail_h: float, img_aspect: float) -> SlotGeometry:
    """
    Dynamically compute left-text / right-image split geometry based on the
    image's ACTUAL aspect ratio.  The image slot adapts to the image so that:
      • Landscape images (wide) get a wider right slot.
      • Portrait images (tall) get a narrower right slot.
      • Text always gets at least 34% of the usable slide width.
      • The image is NEVER forced into a shape that would require cropping.

    Args:
        top:        Top edge of the content zone (EMU).
        avail_h:    Available height for the content zone (EMU).
        img_aspect: width / height of the actual image (clamped internally).

    Returns:
        SlotGeometry with img_* and txt_* fields all in EMU.
    """
    # Clamp to sane bounds (0.4 = very tall portrait, 2.8 = very wide panorama)
    img_aspect = max(0.4, min(2.8, img_aspect))

    # Total usable width: full slide minus left + right margins
    margin_l = Inches(0.25)
    margin_r = Inches(0.3)
    gap      = Inches(0.2)
    total_w  = W - margin_l - margin_r

    # Natural image width if the slot were exactly avail_h tall
    natural_img_w = avail_h * img_aspect

    # Constraints:
    #   • Image can take at most 62% of total_w
    #   • Text must keep at least 34% of total_w
    max_img_w = min(total_w * 0.62, total_w - total_w * 0.34 - gap)
    img_w     = min(natural_img_w, max_img_w)

    # Ensure image is at least 30% of total_w (for very portrait images)
    img_w = max(img_w, total_w * 0.30)

    txt_w = total_w - img_w - gap

    txt_x = margin_l
    img_x = txt_x + txt_w + gap

    return SlotGeometry(
        img_x=img_x, img_y=top, img_w=img_w, img_h=avail_h,
        txt_x=txt_x, txt_y=top, txt_w=txt_w, txt_h=avail_h,
    )

def _parse_bullet(b) -> tuple:
    """Safely parse any bullet item into (bold_text, body_text).
    Handles: proper dicts, single-quoted dict strings, plain strings."""
    if isinstance(b, dict):
        return b.get("bold", ""), b.get("text", b.get("description", ""))
    if isinstance(b, str):
        s = b.strip()
        # Handle Python-style single-quoted dicts: {'bold': '...', 'text': '...'}
        if s.startswith("{") and "'bold'" in s:
            try:
                # Convert single-quote dict repr to valid JSON
                import ast
                d = ast.literal_eval(s)
                if isinstance(d, dict):
                    return d.get("bold", ""), d.get("text", d.get("description", ""))
            except Exception:
                pass
        # Try JSON parse in case it has double quotes
        if s.startswith("{"):
            try:
                d = json.loads(s)
                if isinstance(d, dict):
                    return d.get("bold", ""), d.get("text", d.get("description", ""))
            except Exception:
                pass
        return "", s
    return "", str(b)

def _parse_card(c) -> dict:
    """Safely parse a card item into a dict with 'header' and 'bullets'."""
    if isinstance(c, dict):
        return c
    if isinstance(c, str):
        s = c.strip()
        if s.startswith("{"):
            try:
                import ast
                d = ast.literal_eval(s)
                if isinstance(d, dict): return d
            except Exception:
                pass
            try:
                d = json.loads(s)
                if isinstance(d, dict): return d
            except Exception:
                pass
        return {"header": "Note", "bullets": [s]}
    return {"header": "Note", "bullets": [str(c)]}

def _corner_L(slide, l, t, w, h, col, size=Inches(0.35), th=Inches(0.055)):
    """Draw 4 L-shaped corner brackets."""
    for lx, ly, bw, bh in [
        (l,         t,         size, th),  (l,         t,          th,   size),
        (l+w-size,  t,         size, th),  (l+w-th,    t,          th,   size),
        (l,         t+h-th,    size, th),  (l,         t+h-size,   th,   size),
        (l+w-size,  t+h-th,    size, th),  (l+w-th,    t+h-size,   th,   size),
    ]:
        _rect(slide, lx, ly, bw, bh, fill=col)


# ─── Groq Prompt ───────────────────────────────────────────────────────────────
_SYS_OUTLINE = """\
You are an elite presentation generator. Create a structural outline.
Output ONLY valid JSON — no markdown fences, no explanation."""

_USR_OUTLINE = """\
Create a VISUAL-FIRST presentation outline for:
TOPIC/REQUEST: {prompt}

RULES:
- SLIDE COUNT: {slide_count_rule}
- First slide MUST use "aesthetic_title".
- The FINAL slide (and ONLY the final slide) MUST be a Conclusion/Summary. Do not place it earlier.
- {image_rules}
- Use a mix of layouts: "aesthetic_split", "aesthetic_grid", "aesthetic_flow", "aesthetic_timeline", "aesthetic_comparison", "aesthetic_metrics", "aesthetic_pitch", "aesthetic_poster".
- You MUST use AT LEAST 4 different layouts in the presentation.
- Use "aesthetic_poster" for high-density slides that combine BOTH cards AND bullet points (e.g., problem statement with 4 issue cards + key detail bullets). It automatically builds a dashboard-style layout.
- {purpose_rules}
- DO NOT use the exact same layout for two consecutive slides (e.g., do not put two 'aesthetic_grid' slides back-to-back).
- COLORS: If the user requests ANY specific color or theme (e.g. "cyan", "red", "maroon"), DO NOT rely on presets. You MUST output a "custom_theme" object with 6-character hex codes (NO hash) that PERFECTLY matches their request. Set "ac1", "ac2", "ac3", "border", "hdr_bg", and "bar" to the requested colors.
- Output ONLY valid JSON.

JSON SCHEMA:
{{
  "presentation_title": "...",
  "personality": "ocean_pro",
  "custom_theme": {{"bg": "400000", "bg2": "4A0000", "card": "550000", "card2": "600000", "text": "FFFFFF", "sub": "DDDDDD", "ac1": "FFD700", "ac2": "FFAA00", "ac3": "FFFF00", "border": "FFD700", "hdr_bg": "FFD700", "hdr_text": "400000", "bar": "FFD700"}},
  "slides": [
    {{
      "slide_number": 1,
      "layout": "aesthetic_title",
      "title": "Topic Name",
      "image_slot": false
    }},
    {{
      "slide_number": 2,
      "layout": "aesthetic_split",
      "title": "Problem Statement"
    }}
  ]
}}"""

_SYS_CHUNK = """\
You are an elite presentation generator. Generate highly dense content for specific slides based on the provided outline.
CRITICAL: Do NOT generate the entire presentation. ONLY generate the specific slides requested in the user prompt (from start_idx to end_idx). Returning extra slides is a severe error.
Output ONLY valid JSON — no markdown fences, no explanation."""

_USR_CHUNK = """TOPIC: {prompt}
PRESENTATION OUTLINE: {outline}

Generate the FULL, EXTREMELY DENSE CONTENT for slides {start_idx} to {end_idx}.
Return a JSON object with a "slides" array containing ONLY those slides, fully populated.
The layout for each slide MUST exactly match what is in the outline.

ABSOLUTE RULES:
1. bullets/left_bullets/right_bullets/nodes/cards MUST be JSON arrays of OBJECTS with double-quote keys.
   NEVER output items as plain strings like {{'bold': 'x', 'text': 'y'}} -- always use double quotes.
2. Each bullet: "bold" = 2-5 word sub-header. "text" = 25-35 words of detailed, informative content.
3. Every slide MUST include "visual_suggestion" describing a specific chart, diagram, or infographic.
4. UNIQUE CONTENT: Ensure ALL text is completely unique. Do NOT repeat facts or concepts across different slides OR within the same slide. If a layout requires both cards and bullets, they MUST contain disjoint, complementary information (e.g. cards for high-level summary, bullets for deep-dive technical details).

LAYOUT SCHEMAS:
- aesthetic_showcase: "visual_suggestion": "description of the hero image" (No bullets/cards needed)
- aesthetic_split/aesthetic_pitch/aesthetic_flow: "bullets": [{{"bold":"...", "text":"..."}}] (3-4 bullets)
- aesthetic_grid: "cards": [{{"header":"Card Title", "bullets":["detailed sentence 1", "detailed sentence 2"]}}] (4 cards)
- aesthetic_poster: "cards": [{{"header":"...", "bullets":["..."]}}] (4 cards for broad overview) + "bullets": [{{"bold":"...", "text":"..."}}] (2-3 bullets for deep-dive specific details that DO NOT repeat the cards)
- aesthetic_timeline: "nodes": [{{"header":"Phase Name", "text":"25-35 word description"}}] (4-5 nodes)
- aesthetic_metrics: "metrics": [{{"value":"94%", "label":"Satisfaction Rate"}}] (3-4 metrics)
- aesthetic_comparison: "left_header":"...", "right_header":"...", "left_bullets":[...], "right_bullets":[...]

EXAMPLE OUTPUT for aesthetic_split:
{{"slides": [{{"slide_number": 2, "layout": "aesthetic_split", "title": "Early Life",
  "visual_suggestion": "Timeline diagram with key early life milestones",
  "bullets": [
    {{"bold": "Born 1947 Varanasi", "text": "Pandit Jagdish Mohan was born in the ancient scholarly city of Varanasi in 1947, immersed from childhood in Sanskrit traditions and Vedic philosophy that shaped his entire career."}},
    {{"bold": "University Education", "text": "He studied Sanskrit literature at Banaras Hindu University, earning his degree with distinction and receiving the Chancellor Gold Medal for academic excellence."}},
    {{"bold": "Mentor Influence", "text": "Guided by the legendary Dr. Ramachandra Shukla, he developed rigorous analytical skills in classical Indian philosophy, blending traditional wisdom with modern interdisciplinary methods."}}
  ]
}}]}}"""

_SYS_CHART = """\
You are a data extraction specialist. Given a visual description from a presentation slide,
extract structured numerical/categorical data suitable for chart generation.
Output ONLY valid JSON — no markdown fences, no explanation."""

_USR_CHART = """\
Slide title: {title}
Slide layout: {layout}
Visual suggestion: {visual_suggestion}
Slide content summary: {content_summary}

Extract chart-ready data from the above. Pick the BEST chart type for this data.
If the visual suggestion mentions specific numbers, percentages, or comparisons, extract them.
If there are no clear numerical values, INFER reasonable realistic data from the context.

CRITICAL RULES:
1. NEVER put dates or years (e.g. 1993, 2025) as VALUES in a bar chart (values must be counts, amounts, percentages).
2. VARIETY: You MUST choose a chart type that has NOT been heavily used yet. 
   Already used charts: {used_charts}
   If you see "bar" in the used list, DO NOT use it again. Pick "pie", "line", "timeline", "metrics", etc.
3. LAYOUT MATCHING & REDUNDANCY:
   - For TALL/SQUARE spaces (layout = "aesthetic_split", "aesthetic_pitch"): Use "bar", "pie", or "metrics".
   - For WIDE/HORIZONTAL spaces (layout = "aesthetic_grid", "aesthetic_timeline", "aesthetic_comparison"): Use "line", "timeline", "comparison", or "horizontal_bar". DO NOT use "pie" or "donut" in wide spaces.
   - If the layout is "aesthetic_metrics", the slide ALREADY has large metric text cards. Do NOT generate a "metrics" chart of the same numbers! Generate a complementary "line" or "horizontal_bar" chart instead.
   
4. YOU MUST strictly adhere to the required JSON schemas for your chosen chart type. Failure to do so will result in an empty slide.

Chart types: "bar", "horizontal_bar", "pie", "donut", "line", "comparison", "timeline", "metrics"

RETURN ONLY this JSON:
{{
  "type": "bar",
  "title": "Chart Title",
  "labels": ["Label1", "Label2", "Label3"],
  "values": [94, 85, 78]
}}

For "comparison" type, use:
{{
  "type": "comparison",
  "title": "...",
  "categories": ["Cat1", "Cat2"],
  "left_label": "Group A", "right_label": "Group B",
  "left_values": [90, 80], "right_values": [70, 60]
}}

For "timeline" type, use:
{{
  "type": "timeline",
  "title": "...",
  "nodes": [{{"header": "Phase 1", "text": "Description"}}]
}}

For "metrics" type, use:
{{
  "type": "metrics",
  "title": "...",
  "metrics": [{{"value": "94%", "label": "Satisfaction"}}]
}}"""



def _groq_call(sys_p: str, usr_p: str, tokens: int = 4000) -> str:
    if _GROQ is None: raise RuntimeError("GROQ_API_KEY not configured")
    safe_tokens = min(tokens, 4500)
    try:
        r = _GROQ.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "system", "content": sys_p}, {"role": "user", "content": usr_p}],
            max_tokens=safe_tokens, temperature=0.7, response_format={"type": "json_object"},
        )
        return r.choices[0].message.content or ""
    except Exception as e:
        err_str = str(e).lower()
        if "429" in err_str or "413" in err_str or "rate" in err_str or "too large" in err_str:
            r = _GROQ.chat.completions.create(
                model="llama-3.1-8b-instant",
                messages=[{"role": "system", "content": sys_p}, {"role": "user", "content": usr_p}],
                max_tokens=safe_tokens, temperature=0.7, response_format={"type": "json_object"},
            )
            return r.choices[0].message.content or ""
        raise e

def _parse(raw: str) -> dict:
    raw = re.sub(r"```(?:json)?|```", "", raw).strip()
    return json.loads(raw[raw.find("{"):raw.rfind("}") + 1])


def extract_theme_from_image(image_path: str) -> dict:
    """
    Analyze a PPT screenshot using Groq's vision model and extract
    the exact color palette as a custom_theme dict.
    Returns a dict with keys: bg, bg2, card, card2, text, sub,
    ac1, ac2, ac3, border, hdr_bg, hdr_text, bar
    """
    import base64, pathlib
    if _GROQ is None:
        raise RuntimeError("GROQ_API_KEY not configured")

    img_bytes = pathlib.Path(image_path).read_bytes()
    b64 = base64.standard_b64encode(img_bytes).decode()

    # Detect mime type
    suffix = pathlib.Path(image_path).suffix.lower()
    mime = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
            "webp": "image/webp"}.get(suffix.lstrip("."), "image/png")

    sys_prompt = "You are a color extraction expert. Output ONLY valid JSON — no markdown fences, no explanation."
    usr_prompt = """\
Analyze this presentation slide image and extract its exact color palette.
Identify the following colors from the image:
- bg: The main background color of the slide
- bg2: A slightly darker/lighter secondary background shade
- card: The main card/box background color
- card2: A secondary card color variant
- text: The main body text color
- sub: The subtitle or secondary text color
- ac1: The primary accent color (most prominent — used for headers, borders, highlights)
- ac2: The secondary accent color
- ac3: A tertiary accent color
- border: The border/outline color of cards or boxes
- hdr_bg: The header/title box background color
- hdr_text: The header/title text color
- bar: The bottom bar or footer bar color

Return ONLY this JSON (all values are 6-character hex codes WITHOUT the # symbol):
{
  "bg": "xxxxxx",
  "bg2": "xxxxxx",
  "card": "xxxxxx",
  "card2": "xxxxxx",
  "text": "xxxxxx",
  "sub": "xxxxxx",
  "ac1": "xxxxxx",
  "ac2": "xxxxxx",
  "ac3": "xxxxxx",
  "border": "xxxxxx",
  "hdr_bg": "xxxxxx",
  "hdr_text": "xxxxxx",
  "bar": "xxxxxx"
}"""

    resp = _GROQ.chat.completions.create(
        model="meta-llama/llama-4-scout-17b-16e-instruct",
        messages=[{
            "role": "user",
            "content": [
                {"type": "text", "text": sys_prompt + "\n\n" + usr_prompt},
                {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}}
            ]
        }],
        max_tokens=512,
        temperature=0.1,
    )
    raw = resp.choices[0].message.content or ""
    theme = _parse(raw)
    # Ensure all required keys are present with safe 6-char hex fallbacks
    required = ["bg","bg2","card","card2","text","sub","ac1","ac2","ac3","border","hdr_bg","hdr_text","bar"]
    for k in required:
        v = theme.get(k, "")
        if not isinstance(v, str) or len(v.lstrip("#")) != 6:
            theme[k] = "FFFFFF" if k in ("text","hdr_text") else "111111"
        else:
            theme[k] = v.lstrip("#")
    return theme

# ─────────────────────────────────────────────────────────────────────────────
class PresentationBuilder:
    def __init__(self, plan: dict, out: str):
        self.plan = plan
        self.out = out
        self.prs = Presentation()
        self.prs.slide_width = W
        self.prs.slide_height = H

        # ── Load Style Profile ──
        self.purpose = plan.get("purpose", "hackathon")
        self.S = STYLE_PROFILES.get(self.purpose, STYLE_PROFILES["hackathon"])

        # ── Resolve Palette ──
        pk = plan.get("personality", self.S["default_palette"])
        # If chosen palette doesn't match purpose, switch to default for that purpose
        if pk not in self.S["allowed_palettes"] and not plan.get("custom_theme"):
            pk = self.S["default_palette"]
        base_p = PERSONALITIES.get(pk, PERSONALITIES[self.S["default_palette"]])
        
        # Deep copy to allow overrides
        self.P = dict(base_p)
        
        # Dynamic custom theme overrides! If user specified explicit colors not matching presets
        custom = plan.get("custom_theme", {})
        if isinstance(custom, dict) and custom:
            # Apply ALL colors directly from the LLM's custom_theme.
            # Only fall back for fields the LLM did not provide.
            def _sh(val, fallback):
                """Safe hex: strip # and validate 6 chars, else return fallback."""
                if isinstance(val, str):
                    v = val.lstrip("#")
                    if len(v) == 6:
                        return v
                return fallback

            bg_hex = _sh(custom.get("bg"), "F5F0E8")
            try:
                rv, gv, bv = int(bg_hex[0:2],16), int(bg_hex[2:4],16), int(bg_hex[4:6],16)
                lum = (0.299*rv + 0.587*gv + 0.114*bv) / 255
            except:
                lum = 0.5
            is_light = lum > 0.5

            # Smart defaults based on light/dark detection
            def_text  = "222222" if is_light else "F5F5F5"
            def_sub   = "555555" if is_light else "B0B0C0"
            def_bg2   = "EDE8DC" if is_light else "121218"
            def_card  = "FAF6EE" if is_light else "1A1A22"
            def_card2 = "EDE8DC" if is_light else "22222C"
            def_hdrt  = "222222" if is_light else "FFFFFF"
            def_ac1   = "C8A96E" if is_light else "FFD700"
            def_ac2   = "A07840" if is_light else "FF8800"
            def_ac3   = "8B6530" if is_light else "00C8FF"

            self.P["bg"]       = bg_hex
            self.P["bg2"]      = _sh(custom.get("bg2"),      def_bg2)
            self.P["card"]     = _sh(custom.get("card"),     def_card)
            self.P["card2"]    = _sh(custom.get("card2"),    def_card2)
            self.P["text"]     = _sh(custom.get("text"),     def_text)
            self.P["sub"]      = _sh(custom.get("sub"),      def_sub)
            self.P["hdr_text"] = _sh(custom.get("hdr_text"), def_hdrt)

            ac1 = _sh(custom.get("ac1"), def_ac1)
            ac2 = _sh(custom.get("ac2"), def_ac2)
            ac3 = _sh(custom.get("ac3"), def_ac3)
            self.P["ac1"]    = ac1
            self.P["ac2"]    = ac2
            self.P["ac3"]    = ac3
            self.P["border"] = _sh(custom.get("border"), ac1)
            self.P["hdr_bg"] = _sh(custom.get("hdr_bg"), ac1)
            self.P["bar"]    = _sh(custom.get("bar"),    ac2)
            tc = custom.get("tag_colors", [])
            self.P["tag_colors"] = [
                _sh(tc[0] if len(tc)>0 else None, ac1),
                _sh(tc[1] if len(tc)>1 else None, ac2),
                _sh(tc[2] if len(tc)>2 else None, ac3),
                _sh(tc[3] if len(tc)>3 else None, "888888"),
            ]

            # FORCE HIGH CONTRAST FOR TEXT
            # If the LLM generates a light background but also light text (or dark/dark), fix it.
            def _get_lum(hex_str):
                try:
                    rv, gv, bv = int(hex_str[0:2],16), int(hex_str[2:4],16), int(hex_str[4:6],16)
                    return (0.299*rv + 0.587*gv + 0.114*bv) / 255
                except:
                    return 0.5
            
            text_lum = _get_lum(self.P["text"])
            # If contrast between background and text is too low, force safe defaults
            if abs(text_lum - lum) < 0.4:
                self.P["text"] = "111111" if is_light else "F5F5F5"
                self.P["sub"]  = "444444" if is_light else "B0B0C0"
                self.P["hdr_text"] = "111111" if is_light else "FFFFFF"

    def _blank(self):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        _bg_fill(s, self.P["bg"])
        if getattr(self, "S", {}).get("gradient_strips", True):
            _rect(s, 0, 0, W, Inches(0.06), fill=self.P["ac1"])
            _rect(s, 0, H - Inches(0.28), W, Inches(0.28), fill=self.P["bar"])
        else:
            _rect(s, 0, H - Inches(0.04), W, Inches(0.04), fill=self.P["ac1"])
            
        if getattr(self, "S", {}).get("right_edge_panel", True):
            _rect(s, W - Inches(0.18), Inches(0.06), Inches(0.18), H - Inches(0.34), fill=self.P["bg2"])
        return s

    def _progress(self, slide, cur, total):
        _tb(slide, f"{cur:02d} / {total:02d}", W - Inches(1.35), H - Inches(0.24),
            Inches(1.1), Inches(0.2), sz=9, bold=True, col=self.P["text"], align=PP_ALIGN.RIGHT)

    def _header(self, slide, title: str):
        if getattr(self, "S", {}).get("header_style", "pill") == "pill":
            tw, th = W * 0.65, Inches(0.65)
            tx, ty = (W - tw) / 2, Inches(0.15)
            _round(slide, tx, ty, tw, th, fill=self.P["card"], line=self.P["border"], lw=1.5)
            _tb(slide, title.upper(), tx, ty + Inches(0.08), tw, th - Inches(0.16),
                sz=22, bold=True, col=self.P["text"], align=PP_ALIGN.CENTER)
        else:
            bw = getattr(self, "S", {}).get("card_border_width", 0.75)
            sz = getattr(self, "S", {}).get("header_font_size", 20)
            if getattr(self, "S", {}).get("filled_header", True):
                _rect(slide, Inches(0.25), Inches(0.15), W - Inches(0.5), Inches(0.55), fill=self.P["hdr_bg"], line=self.P["border"], lw=bw)
                _tb(slide, title.upper(), Inches(0.45), Inches(0.22), W - Inches(0.9), Inches(0.4), sz=sz, bold=True, col=self.P["hdr_text"], align=PP_ALIGN.LEFT)
            else:
                # Clean minimal underline
                _rect(slide, Inches(0.25), Inches(0.7), W - Inches(0.5), Inches(0.02), fill=self.P["ac1"])
                _tb(slide, title.upper(), Inches(0.25), Inches(0.15), W - Inches(0.5), Inches(0.5), sz=sz, bold=True, col=self.P["text"], align=PP_ALIGN.LEFT)

    def _premium_image_frame(self, slide, path_or_paths, l, t, w, h, suggestion: str, chart_data: dict = None):
        """
        Draws one or multiple premium images intelligently tiled within the bounding box.
        """
        if isinstance(path_or_paths, str):
            paths = [path_or_paths]
        else:
            paths = path_or_paths
            
        # Filter and clean
        paths = [_clean_image_path(p) for p in paths if p]
        
        if not paths:
            self._premium_image_frame_single(slide, "", l, t, w, h, suggestion, chart_data)
            return
            
        n = len(paths)
        if n == 1:
            self._premium_image_frame_single(slide, paths[0], l, t, w, h, suggestion, chart_data)
            return
            
        # ── SMART MASONRY TILING ──────────────────────────────────────────────
        gap = Inches(0.18)
        
        # Calculate aspect ratios for all images (fallback to 16:9 if invalid)
        aspects = []
        for p in paths:
            r = _get_image_aspect_ratio(p)
            if r <= 0: r = 1.78
            aspects.append(r)
            
        # Option A: Horizontal Gallery (All images have the SAME HEIGHT)
        h_a = (w - gap * (n - 1)) / sum(aspects)
        h_a = min(h_a, h)
        w_a = sum([h_a * r for r in aspects]) + gap * (n - 1)
        area_a = sum([h_a * (h_a * r) for r in aspects])
        
        # Option B: Vertical Gallery (All images have the SAME WIDTH)
        w_b = (h - gap * (n - 1)) / sum([1.0/r for r in aspects])
        w_b = min(w_b, w)
        h_b = sum([w_b / r for r in aspects]) + gap * (n - 1)
        area_b = sum([w_b * (w_b / r) for r in aspects])
        
        if area_a >= area_b:
            # Render Option A (Horizontal, perfectly centered)
            cx = l + (w - w_a) / 2
            cy = t + (h - h_a) / 2
            curr_x = cx
            for i, p in enumerate(paths):
                img_w = h_a * aspects[i]
                self._premium_image_frame_single(slide, p, curr_x, cy, img_w, h_a, 
                                                 suggestion if i == 0 else "", 
                                                 chart_data if i == 0 else None)
                curr_x += img_w + gap
        else:
            # Render Option B (Vertical, perfectly centered)
            cx = l + (w - w_b) / 2
            cy = t + (h - h_b) / 2
            curr_y = cy
            for i, p in enumerate(paths):
                img_h = w_b / aspects[i]
                self._premium_image_frame_single(slide, p, cx, curr_y, w_b, img_h, 
                                                 suggestion if i == 0 else "", 
                                                 chart_data if i == 0 else None)
                curr_y += img_h + gap

    def _premium_image_frame_single(self, slide, path: str, l, t, w, h, suggestion: str, chart_data: dict = None):
        """
        Draws a single premium rounded frame. Images are perfectly shrink-wrapped.
        """
        bw = getattr(self, "S", {}).get("card_border_width", 1.5)
        frame_pad = Inches(0.12)      # uniform inner padding around the frame border

        clean = _clean_image_path(path) if path else ""
        if clean and Path(clean).exists():
            try:
                # ── Step 1: Get real image dimensions ──
                try:
                    from PIL import Image as _PIL
                    with _PIL.open(clean) as _im:
                        img_px_w, img_px_h = _im.size
                except Exception:
                    img_px_w, img_px_h = 16, 9   # safe 16:9 assumption

                img_ratio  = img_px_w / max(img_px_h, 1)
                
                # Dynamically cap padding so it never exceeds 10% of the image size
                # This prevents negative geometry bugs when drawing thin/small images
                actual_pad = min(frame_pad, w * 0.1, h * 0.1)
                
                inner_w    = max(w - 2 * actual_pad, 0.01)
                inner_h    = max(h - 2 * actual_pad, 0.01)
                slot_ratio = inner_w / inner_h

                # ── Step 2: Contain fit to find optimal dimensions ──
                if img_ratio >= slot_ratio:
                    fit_w = inner_w
                    fit_h = inner_w / img_ratio
                else:
                    fit_h = inner_h
                    fit_w = inner_h * img_ratio

                # ── Step 3: Shrink-wrap the frame tightly around the image ──
                final_w = fit_w
                final_h = fit_h
                
                # Center the tight frame within the original (l, t, w, h) slot
                frame_x = l + (w - final_w) / 2
                frame_y = t + (h - final_h) / 2

                # Insert the image perfectly filling the inner frame
                from pptx.enum.shapes import MSO_SHAPE
                pic = slide.shapes.add_picture(
                    clean,
                    frame_x,
                    frame_y,
                    final_w,
                    final_h,
                )
                
                # Give the image itself rounded corners so it doesn't bleed out!
                pic.auto_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
                
                # Apply premium border directly to the image
                pic.line.color.rgb = _c(self.P["border"])
                pic.line.width = Pt(bw)

                if getattr(self, "S", {}).get("tech_dots", True):
                    for di in range(3):
                        _oval(slide, frame_x + final_w - Inches(0.6) + di * Inches(0.16), frame_y + Inches(0.15),
                              Inches(0.08), Inches(0.08), fill=[self.P["ac1"], self.P["ac2"], self.P["ac3"]][di], line=None)

                return
            except Exception as e:
                print(f"[ppt] image error: {e}")

        # ── AUTO-CHART / FALLBACK (Uses full box) ──
        _round(slide, l, t, w, h, fill=self.P["card"], line=self.P["border"], lw=bw)
        if getattr(self, "S", {}).get("tech_dots", True):
            for di in range(3):
                _oval(slide, l + w - Inches(0.6) + di * Inches(0.16), t + Inches(0.15),
                      Inches(0.08), Inches(0.08), fill=[self.P["ac1"], self.P["ac2"], self.P["ac3"]][di], line=None)

        if chart_data and isinstance(chart_data, dict):
            try:
                from app.services.ppt_chart_engine import ChartEngine
                png_bytes = ChartEngine.render(
                    chart_data, 
                    self.P, 
                    w=(w / 914400.0) - 0.5, 
                    h=(h / 914400.0) - 0.5
                )
                if png_bytes and len(png_bytes) > 100:
                    c_pad = Inches(0.25)
                    pic = slide.shapes.add_picture(
                        io.BytesIO(png_bytes), l + c_pad, t + c_pad,
                        w - 2*c_pad, h - 2*c_pad
                    )
                    return
            except Exception as e:
                print(f"[ppt] chart render error: {e}")

        # Text fallback
        try:
            from pptx.enum.text import MSO_ANCHOR
            v_align = MSO_ANCHOR.MIDDLE
        except ImportError:
            v_align = None
            
        _tb(slide, suggestion or "[ Detailed Visual ]",
            l + Inches(0.2), t + Inches(0.2), w - Inches(0.4), h - Inches(0.4),
            sz=14, italic=False, col=self.P["sub"], align=PP_ALIGN.CENTER, v_align=v_align)

    def _bullet_card(self, slide, l, t, w, h, bold_txt: str, body_txt: str, tag_col: str, idx: int):
        sz_b = getattr(self, "S", {}).get("body_font_size", 11)
        
        # Ensure text is visible against card background
        def _lum(c):
            try: return (0.299*int(c[0:2],16) + 0.587*int(c[2:4],16) + 0.114*int(c[4:6],16))/255
            except: return 0.5
        card_col = self.P.get("card", "FFFFFF")
        text_col = self.P.get("text", "000000")
        if abs(_lum(text_col) - _lum(card_col)) < 0.35:
            text_col = "111111" if _lum(card_col) > 0.5 else "F5F5F5"

        if getattr(self, "S", {}).get("bullet_style", "card") == "card":
            _round(slide, l, t, w, h, fill=self.P["card"], line=tag_col, lw=1.5)
            if getattr(self, "S", {}).get("numbered_badges", True):
                _oval(slide, l + Inches(0.15), t + Inches(0.12), Inches(0.25), Inches(0.25), fill=tag_col, line=None)
                _tb(slide, str(idx), l + Inches(0.15), t + Inches(0.12), Inches(0.25), Inches(0.23), sz=10, bold=True, col=self.P["bg"], align=PP_ALIGN.CENTER)
            if bold_txt:
                _tb(slide, bold_txt, l + Inches(0.5), t + Inches(0.1), w - Inches(0.6), Inches(0.3), sz=12, bold=True, col=tag_col)
            body_top = t + Inches(0.35) if bold_txt else t + Inches(0.14)
            _tb(slide, body_txt, l + Inches(0.5), body_top, w - Inches(0.6), h - (Inches(0.4) if bold_txt else Inches(0.22)), sz=sz_b, col=text_col)
        else:
            if getattr(self, "S", {}).get("numbered_badges", False):
                _tb(slide, f"{idx}.", l, t + Inches(0.05), Inches(0.3), Inches(0.25), sz=12, bold=True, col=tag_col, align=PP_ALIGN.RIGHT)
            else:
                _oval(slide, l + Inches(0.15), t + Inches(0.1), Inches(0.1), Inches(0.1), fill=tag_col, line=None)
            tx, tw = l + Inches(0.35), w - Inches(0.4)
            if bold_txt:
                _tb(slide, bold_txt, tx, t, tw, Inches(0.25), sz=12, bold=True, col=tag_col)
                _tb(slide, body_txt, tx, t + Inches(0.25), tw, h - Inches(0.25), sz=sz_b, col=self.P["text"])
            else:
                _tb(slide, body_txt, tx, t, tw, h, sz=sz_b, col=self.P["text"])

    def build_with_progress(self) -> Generator[str, None, None]:
        slides = self.plan.get("slides", [])
        total = len(slides)

        # ── Render each slide ──────────────────────────────────────────────────
        for i, sd in enumerate(slides):
            lay = sd.get("layout", "aesthetic_split")
            yield f"Slide {i+1} - {sd.get('title', '')}"

            fn = getattr(self, f"_lay_{lay}", self._lay_aesthetic_split)
            try:
                fn(sd, i+1, total)
            except Exception as e:
                yield f"Error {lay}: {e}"
                try: self._lay_aesthetic_split(sd, i+1, total)
                except: pass
        yield "Saving..."
        self.prs.save(self.out)

    # ── LAYOUT 1: TITLE SLIDE ─────────────────────────────────────────────────
    def _lay_aesthetic_title(self, d: dict, cur: int, total: int):
        slide = self._blank()
        self._progress(slide, cur, total)

        has_img = bool(d.get("image_path") or d.get("image_slot"))

        if has_img:
            # Side-by-side title slide with image
            img_w = W * 0.45
            txt_w = W - img_w - Inches(0.8)
            cw, ch = txt_w, Inches(5.4)
            cx, cy = Inches(0.4), (H - ch) / 2
            
            # Render image on the right
            self._premium_image_frame(
                slide, d.get("image_path", ""), W - img_w - Inches(0.4), Inches(0.8),
                img_w, H - Inches(1.6), d.get("visual_suggestion", "Hero Image")
            )
        else:
            # Centered title slide
            cw, ch = Inches(11.0), Inches(5.4)
            cx, cy = (W - cw)/2, (H - ch)/2 - Inches(0.1)

            if getattr(self, "S", {}).get("decorative_circles", True):
                _oval(slide, W*0.55, H*0.1, Inches(5.5), Inches(5.5), fill=self.P["bg2"], line=self.P["border"], lw=0.5)
                _oval(slide, W*0.58, H*0.18, Inches(4.0), Inches(4.0), fill=self.P["bg"], line=self.P["ac2"], lw=0.4)

        if getattr(self, "S", {}).get("drop_shadows", True):
            _round(slide, cx + Inches(0.15), cy + Inches(0.15), cw, ch, fill=self.P["bg2"], line=None)
            
        if getattr(self, "S", {}).get("heavy_title_card", True):
            bw = getattr(self, "S", {}).get("card_border_width", 1.5)
            _round(slide, cx, cy, cw, ch, fill=self.P["card"], line=self.P["border"], lw=bw)
            _round(slide, cx, cy, cw, Inches(0.2), fill=self.P["ac1"])
            _rect(slide, cx, cy + ch - Inches(0.08), cw, Inches(0.08), fill=self.P["ac2"])
        else:
            # Minimal general title
            _rect(slide, cx + Inches(2.0), cy + ch/2 - Inches(0.4), cw - Inches(4.0), Inches(0.02), fill=self.P["ac1"])

        if getattr(self, "S", {}).get("corner_brackets", True):
            _corner_L(slide, cx, cy, cw, ch, col=self.P["ac1"], size=Inches(0.5), th=Inches(0.07))

        _tb(slide, d.get("title", "Presentation"), cx + Inches(0.5), cy + Inches(0.4),
            cw - Inches(1.0), Inches(2.4), sz=getattr(self, "S", {}).get("title_font_size", 48) if not has_img else 40, bold=True, col=self.P["text"],
            align=PP_ALIGN.CENTER if not has_img else PP_ALIGN.LEFT, font=getattr(self, "S", {}).get("title_font", "Calibri"))

        if not has_img:
            _rect(slide, cx + Inches(2.2), cy + Inches(2.95), cw - Inches(4.4), Inches(0.05), fill=self.P["ac1"])
        
        _tb(slide, d.get("subtitle", ""), cx + Inches(0.7), cy + Inches(3.1),
            cw - Inches(1.4), Inches(1.95), sz=13.5, col=self.P["sub"], align=PP_ALIGN.CENTER if not has_img else PP_ALIGN.LEFT)

        if getattr(self, "S", {}).get("tech_dots", True):
            for di in range(3):
                _rect(slide, cx + cw - Inches(0.8) + di * Inches(0.22),
                      cy + ch - Inches(0.38), Inches(0.12), Inches(0.12), fill=self.P["ac1"])

    # ── LAYOUT 2: SHOWCASE — MASSIVE HERO IMAGE (or side-by-side pair) ───────
    def _lay_aesthetic_showcase(self, d: dict, cur: int, total: int):
        slide = self._blank()
        self._header(slide, d.get("title", ""))
        self._progress(slide, cur, total)

        top     = Inches(0.86)
        avail_h = H - top - Inches(0.3)
        margin  = Inches(0.25)
        gap     = Inches(0.18)
        avail_w = W - 2 * margin

        extra_paths = d.get("extra_image_paths", [])
        all_paths = [d.get("image_path", "")] + list(extra_paths) if d.get("image_path") or extra_paths else []
        
        self._premium_image_frame(
            slide, all_paths,
            margin, top,
            avail_w, avail_h,
            d.get("visual_suggestion", "[ Massive Hero Visual ]")
        )


    # ── LAYOUT: SPLIT — INTELLIGENT TEXT + IMAGE ────────────────────────────
    def _lay_aesthetic_split(self, d: dict, cur: int, total: int):
        """
        Fluid split layout that adapts natively to image aspect ratio:
        • Landscape (>1.2): Image spans top, text in 2 columns below.
        • Portrait/Square (<=1.2): Image on left/right, text beside it.
        """
        bullets = [b for b in d.get("bullets", []) if b]
        
        # ── INTELLIGENT CONTENT FALLBACK ──
        if not bullets and d.get("description"):
            bullets = [{"bold": "Overview", "text": d.get("description")}]
        elif not bullets and d.get("subtitle"):
            bullets = [{"bold": "Summary", "text": d.get("subtitle")}]
            
        if not bullets or all(not _parse_bullet(b)[1] for b in bullets):
            self._lay_aesthetic_showcase(d, cur, total)
            return

        slide = self._blank()
        self._header(slide, d.get("title", ""))
        self._progress(slide, cur, total)

        top   = Inches(0.86)
        avail = H - top - Inches(0.3)
        margin = Inches(0.25)
        avail_w = W - 2 * margin
        gap = Inches(0.15)

        img_path   = _clean_image_path(d.get("image_path", ""))
        has_image  = bool(img_path and Path(img_path).exists())
        img_aspect = float(d.get("image_aspect", 0.0))
        if img_aspect <= 0 and has_image:
            img_aspect = _get_image_aspect_ratio(img_path)
        if img_aspect <= 0:
            img_aspect = 1.78

        tag_colors = self.P["tag_colors"]

        # ALWAYS use Left/Right split for standard slides. Vertical mode squishes the cards too much.
        geo = compute_split_geometry(top, avail, img_aspect)
        
        # Alternate sides
        if cur % 2 == 0:
            img_x = margin
            txt_x = margin + geo.img_w + gap
        else:
            txt_x = geo.txt_x
            img_x = geo.img_x

        all_paths = [d.get("image_path", "")] + d.get("extra_image_paths", []) if has_image else []
        
        # Image frame fills full content height — _premium_image_frame handles
        # letterbox containment internally, so no wasted white space.
        self._premium_image_frame(
            slide, all_paths,
            img_x, top, geo.img_w, avail,
            d.get("visual_suggestion", ""),
            chart_data=d.get("chart_data"),
        )

        # Bullet cards
        n = max(len(bullets[:4]), 1)
        ch2 = (geo.txt_h - gap * (n - 1)) / n
        by = geo.txt_y

        for i, b in enumerate(bullets[:4]):
            bold_txt, body_txt = _parse_bullet(b)
            self._bullet_card(
                slide, txt_x, by, geo.txt_w, ch2,
                bold_txt, body_txt,
                tag_colors[i % len(tag_colors)], i + 1,
            )
            by += ch2 + gap

    # ── LAYOUT 3: GRID — 2×2 COLORED CARDS (+ optional side image) ───────────
    def _lay_aesthetic_grid(self, d: dict, cur: int, total: int):
        slide = self._blank()
        self._header(slide, d.get("title", ""))
        self._progress(slide, cur, total)

        top   = Inches(0.86)
        avail = H - top - Inches(0.3)
        margin = Inches(0.25)
        gap   = Inches(0.18)
        cards = d.get("cards", [])
        
        if not cards and "bullets" in d:
            for i, b in enumerate(d["bullets"]):
                bold, text = _parse_bullet(b)
                header = bold if bold else (f"Section {i+1}")
                cards.append({"header": header, "bullets": [text] if text else ["(No content)"]})
        
        cards = [_parse_card(c) for c in cards]
        for i, c in enumerate(cards):
            if not c.get("header", ""):
                c["header"] = f"Module {i+1}"

        # Check if we have an image to show alongside
        img_path  = _clean_image_path(d.get("image_path", ""))
        has_img   = bool(img_path and Path(img_path).exists())
        has_chart = bool(d.get("chart_data"))

        if has_img or has_chart:
            img_aspect = float(d.get("image_aspect", 0.0))
            if img_aspect <= 0 and has_img:
                img_aspect = _get_image_aspect_ratio(img_path)
            if img_aspect <= 0:
                img_aspect = 1.78
                
            geo = compute_split_geometry(top, avail, img_aspect)
            
            if cur % 2 != 0:
                img_x = margin
                txt_x = margin + geo.img_w + gap
            else:
                txt_x = geo.txt_x
                img_x = geo.img_x
                
            all_paths = [d.get("image_path", "")] + d.get("extra_image_paths", []) if has_img else []
            
            # Image frame fills full content height — _premium_image_frame handles letterbox
            self._premium_image_frame(
                slide, all_paths,
                img_x, top, geo.img_w, avail,
                d.get("visual_suggestion", ""),
                chart_data=d.get("chart_data")
            )
            
            n_cols, n_rows = 2, 2
            col_w = (geo.txt_w - gap) / 2
            row_h = (geo.txt_h - gap) / 2
            
            for i, c in enumerate(cards[:4]):
                ci, ri = i % n_cols, i // n_cols
                cx2 = txt_x + ci * (col_w + gap)
                cy2 = geo.txt_y + ri * (row_h + gap)
                self._colored_card_full(slide, cx2, cy2, col_w, row_h, c, self.P["tag_colors"][i % 4])
        else:
            # Standard 2×2 grid (no chart data, no image)
            n_cols, n_rows = 2, 2
            col_w = (W - Inches(0.5) - gap) / 2
            row_h = (avail - gap) / 2
            for i, c in enumerate(cards[:4]):
                ci, ri = i % n_cols, i // n_cols
                cx2 = Inches(0.25) + ci * (col_w + gap)
                cy2 = top + ri * (row_h + gap)
                self._colored_card_full(slide, cx2, cy2, col_w, row_h, c,
                                        self.P["tag_colors"][i % 4])

    # ── LAYOUT: POSTER — Dynamic Dashboard (Cards + Image + Bullets together) ──
    def _lay_aesthetic_poster(self, d: dict, cur: int, total: int):
        """
        Dynamic constraint-based layout engine.

        Analyses the slide content at runtime and divides the canvas into
        logical regions to fit everything without overlap or wasted space:

          • If cards are present  → horizontal card band (top strip, 1×N)
          • Remaining space below → dynamic left/right split
              – Left:  bullet point cards, stacked vertically
              – Right: image, sized to its REAL aspect ratio
          • If no image: full-width cards across two rows
          • If no cards: falls back to aesthetic_split

        All geometry is computed from the actual image aspect ratio and
        content count — nothing is hardcoded.
        """
        slide = self._blank()
        self._header(slide, d.get("title", ""))
        self._progress(slide, cur, total)

        MX   = Inches(0.25)       # left margin
        MXR  = Inches(0.2)        # right margin
        TOP  = Inches(0.86)       # below header
        BOT  = Inches(0.3)        # above footer bar
        GAP  = Inches(0.15)       # gap between elements
        FULL_W = W - MX - MXR    # usable width

        avail_h = H - TOP - BOT

        bullets = [b for b in d.get("bullets", []) if b]
        
        # ── INTELLIGENT CONTENT FALLBACK ──
        # If the LLM didn't generate bullets for the bottom section, dynamically pull the 
        # description or subtitle so we can populate a summary card next to the image.
        # This prevents the image from floating in a massive empty void.
        if not bullets and d.get("description"):
            bullets = [{"bold": "Overview", "text": d.get("description")}]
        elif not bullets and d.get("subtitle"):
            bullets = [{"bold": "Summary", "text": d.get("subtitle")}]
            
        cards   = [_parse_card(c) for c in d.get("cards", []) if c]
        for i, c in enumerate(cards):
            if not c.get("header"):
                c["header"] = f"Module {i+1}"

        img_path  = _clean_image_path(d.get("image_path", ""))
        has_img   = bool(img_path and Path(img_path).exists())
        has_chart = bool(d.get("chart_data"))
        tag_colors = self.P["tag_colors"]

        # ── No cards + no image → fall back to text-only aesthetic_pitch ──────
        if not cards and not has_img and not has_chart:
            self._lay_aesthetic_pitch(d, cur, total)
            return

        # ── No bullets + no cards → full-screen image showcase ───────────────
        if not bullets and not cards and (has_img or has_chart):
            self._lay_aesthetic_showcase(d, cur, total)
            return

        # ─────────────────────────────────────────────────────────────────────
        # REGION PLANNING
        # We always work top→bottom. The card band (if cards exist) goes first,
        # then the remaining space is split left (bullets) / right (image).
        # ─────────────────────────────────────────────────────────────────────

        # Decide how many cards to show and how many columns to use in the band
        n_cards = min(len(cards), 4)

        if cards and n_cards > 0:
            # Height of the card band is 30–38% of avail_h depending on content density
            card_band_h = avail_h * 0.33 if bullets or has_img else avail_h
            card_y      = TOP
            content_y   = TOP + card_band_h + GAP
            content_h   = avail_h - card_band_h - GAP
        else:
            card_band_h = 0
            card_y      = TOP
            content_y   = TOP
            content_h   = avail_h

        # ── Draw card band ────────────────────────────────────────────────────
        if n_cards > 0:
            col_gap   = GAP
            card_cols = n_cards  # 1 card per column in the top band (max 4)
            card_w    = (FULL_W - col_gap * (card_cols - 1)) / card_cols
            for i, c in enumerate(cards[:n_cards]):
                cx = MX + i * (card_w + col_gap)
                self._colored_card_full(
                    slide, cx, card_y, card_w, card_band_h,
                    c, tag_colors[i % len(tag_colors)]
                )

        # ── Content zone (bullets + image) below the card band ────────────────
        if has_img or has_chart:
            # Compute image aspect ratio
            img_aspect = float(d.get("image_aspect", 0.0))
            if img_aspect <= 0 and has_img:
                img_aspect = _get_image_aspect_ratio(img_path)
            if img_aspect <= 0:
                img_aspect = 1.78

            # How wide should the image be?
            if not bullets:
                img_w = FULL_W
                txt_w = 0
            else:
                nat_img_w = content_h * img_aspect
                # Cap at 62% of full width, ensure at least 38% for bullets
                max_img_w = min(FULL_W * 0.62, FULL_W - FULL_W * 0.34 - GAP)
                img_w = min(nat_img_w, max_img_w)
                img_w = max(img_w, FULL_W * 0.35)  # at least 35%
                txt_w = FULL_W - img_w - GAP

            # Actual image height to preserve aspect ratio (never stretch)
            actual_img_h = img_w / img_aspect
            # Centre image vertically in the content zone
            img_y = content_y + (content_h - actual_img_h) / 2

            # Alternate left/right per slide for visual variety
            if cur % 2 == 0:
                img_x = MX
                txt_x = MX + img_w + GAP if img_w < FULL_W else MX
            else:
                txt_x = MX
                img_x = MX + txt_w + GAP if img_w < FULL_W else MX

            all_paths = [d.get("image_path", "")] + d.get("extra_image_paths", []) if has_img else []
            
            # Image frame fills the full content zone height — letterbox handled internally
            self._premium_image_frame(
                slide, all_paths,
                img_x, content_y, img_w, content_h,
                d.get("visual_suggestion", ""),
                chart_data=d.get("chart_data"),
            )

            # Draw bullet cards in text zone
            if bullets:
                n_b  = min(len(bullets), 4)
                
                # Cap the maximum height of a bullet card so sparse text doesn't look awkwardly huge
                max_bh = Inches(1.5)
                bh   = min((content_h - GAP * (n_b - 1)) / max(n_b, 1), max_bh)
                
                # Vertically center the entire stack of bullets so they sit elegantly next to the image
                total_bullets_h = n_b * bh + (n_b - 1) * GAP
                start_y = content_y + (content_h - total_bullets_h) / 2
                
                by   = start_y
                for i, b in enumerate(bullets[:n_b]):
                    bold_txt, body_txt = _parse_bullet(b)
                    self._bullet_card(
                        slide, txt_x, by, txt_w, bh,
                        bold_txt, body_txt,
                        tag_colors[i % len(tag_colors)], i + 1
                    )
                    by += bh + GAP
            elif not cards:
                # No bullets → show description in text zone
                desc = d.get("description", d.get("subtitle", ""))
                if desc:
                    _round(slide, txt_x, content_y, txt_w, content_h,
                           fill=self.P["card"], line=self.P["border"], lw=1.5)
                    _tb(slide, desc, txt_x + Inches(0.2), content_y + Inches(0.2),
                        txt_w - Inches(0.4), content_h - Inches(0.4),
                        sz=11, col=self.P["text"])

        else:
            # No image → full-width bullet grid below the card band
            if bullets:
                n_b   = min(len(bullets), 6)
                # Use 2 columns if more than 3 bullets for better use of space
                cols  = 2 if n_b > 3 else 1
                rows  = math.ceil(n_b / cols)
                bw    = (FULL_W - GAP * (cols - 1)) / cols
                bh    = (content_h - GAP * (rows - 1)) / max(rows, 1)
                for i, b in enumerate(bullets[:n_b]):
                    col_i = i % cols
                    row_i = i // cols
                    bx    = MX + col_i * (bw + GAP)
                    by    = content_y + row_i * (bh + GAP)
                    bold_txt, body_txt = _parse_bullet(b)
                    self._bullet_card(
                        slide, bx, by, bw, bh,
                        bold_txt, body_txt,
                        tag_colors[i % len(tag_colors)], i + 1
                    )


    def _colored_card_full(self, slide, l, t, w, h, card: dict, color: str):
        """Clean full grid card without messy overlapping header bands."""
        
        # Ensure text is visible against card background
        def _lum(c):
            try: return (0.299*int(c[0:2],16) + 0.587*int(c[2:4],16) + 0.114*int(c[4:6],16))/255
            except: return 0.5
        card_col = self.P.get("card", "FFFFFF")
        text_col = self.P.get("text", "000000")
        if abs(_lum(text_col) - _lum(card_col)) < 0.35:
            text_col = "111111" if _lum(card_col) > 0.5 else "F5F5F5"
            
        # Also ensure header color is visible against card background
        if abs(_lum(color) - _lum(card_col)) < 0.25:
            color = "111111" if _lum(card_col) > 0.5 else "FFFFFF"

        bw = getattr(self, "S", {}).get("card_border_width", 1.5)
        border_col = color if getattr(self, "S", {}).get("colored_borders", True) else self.P["border"]
        _round(slide, l, t, w, h, fill=self.P["card"], line=border_col, lw=bw)

        hh = Inches(0.4)
        _tb(slide, card.get("header", "Module").upper(), l + Inches(0.15), t + Inches(0.1),
            w - Inches(0.3), hh, sz=12, bold=True, col=color, align=PP_ALIGN.LEFT)
            
        _rect(slide, l + Inches(0.15), t + Inches(0.42), w - Inches(0.3), Inches(0.02), fill=color)

        if getattr(self, "S", {}).get("tech_dots", True):
            _oval(slide, l + w - Inches(0.3), t + Inches(0.15), Inches(0.15), Inches(0.15), fill=color, line=None)

        bullets = card.get("bullets", [])
        if not bullets and "text" in card:
            bullets = [card["text"]]
        elif not bullets and "description" in card:
            bullets = [card["description"]]
        elif not bullets:
            bullets = ["(No detailed content provided)"]
            
        avail_h = h - hh - Inches(0.15)
        bsh = avail_h / max(len(bullets), 1)
        by = t + hh + Inches(0.1)
        
        sz_b = getattr(self, "S", {}).get("body_font_size", 11.5)
        
        # Combine bullets into a single paragraph block to prevent overlap
        combined = []
        is_numbered = getattr(self, "S", {}).get("numbered_badges", True)
        
        for bi, b in enumerate(bullets[:5]):
            bold_part, text_part = _parse_bullet(b)
            display = text_part if text_part else (bold_part if bold_part else str(b))
            prefix = f"{bi+1}. " if is_numbered else "• "
            combined.append(f"{prefix}{display}")
            
        full_text = "\n\n".join(combined)
        
        # Dynamically shrink font size if text is exceptionally long to ensure fit
        if len(full_text) > 250:
            sz_b = max(sz_b - 2, 8)
            
        # Draw a single auto-fitting text box to handle line wrapping naturally
        _tb(slide, full_text, l + Inches(0.15), by, w - Inches(0.3), avail_h, sz=sz_b, col=text_col)
    # ── LAYOUT 4: FLOW — TOP CALLOUT + DOMINANT FULL-WIDTH IMAGE ─────────────
    def _lay_aesthetic_flow(self, d: dict, cur: int, total: int):
        slide = self._blank()
        self._header(slide, d.get("title", ""))
        self._progress(slide, cur, total)

        top   = Inches(0.86)
        avail = H - top - Inches(0.3)

        desc_h = Inches(1.5)
        if getattr(self, "S", {}).get("drop_shadows", True):
            _round(slide, Inches(0.25) + Inches(0.07), top + Inches(0.07), W - Inches(0.43), desc_h, fill=self.P["bg2"], line=None)
            
        bw = getattr(self, "S", {}).get("card_border_width", 1.8)
        _round(slide, Inches(0.25), top, W - Inches(0.43), desc_h, fill=self.P["card"], line=self.P["border"], lw=bw)
        
        if getattr(self, "S", {}).get("gradient_strips", True):
            _rect(slide, Inches(0.25), top + Inches(0.06), Inches(0.1), desc_h - Inches(0.12), fill=self.P["ac2"])
            
        _round(slide, Inches(0.44), top + Inches(0.1), Inches(0.35), Inches(0.35), fill=self.P["ac2"], line=None)
        _tb(slide, '"', Inches(0.46), top + Inches(0.09), Inches(0.33), Inches(0.33), sz=18, bold=True, col=self.P["bg"], align=PP_ALIGN.CENTER)
        desc_text = d.get("description", "")
        # Fallback if LLM generated "text" or "bullets" instead of "description"
        if not desc_text:
            desc_text = d.get("text", "")
        if not desc_text and "bullets" in d:
            parts = []
            for b in d["bullets"]:
                _, txt = _parse_bullet(b)
                parts.append(txt)
            desc_text = " ".join(parts)
        if not desc_text and "nodes" in d:
            desc_text = " ".join([n.get("text", "") if isinstance(n, dict) else str(n) for n in d["nodes"]])

        if desc_text:
            if getattr(self, "S", {}).get("drop_shadows", True):
                _round(slide, Inches(0.25) + Inches(0.07), top + Inches(0.07), W - Inches(0.43), desc_h, fill=self.P["bg2"], line=None)

            bw = getattr(self, "S", {}).get("card_border_width", 1.8)
            _round(slide, Inches(0.25), top, W - Inches(0.43), desc_h, fill=self.P["card"], line=self.P["border"], lw=bw)

            if getattr(self, "S", {}).get("gradient_strips", True):
                _rect(slide, Inches(0.25), top + Inches(0.06), Inches(0.1), desc_h - Inches(0.12), fill=self.P["ac2"])

            _round(slide, Inches(0.44), top + Inches(0.1), Inches(0.35), Inches(0.35), fill=self.P["ac2"], line=None)
            _tb(slide, '"', Inches(0.46), top + Inches(0.09), Inches(0.33), Inches(0.33), sz=18, bold=True, col=self.P["bg"], align=PP_ALIGN.CENTER)
            _tb(slide, desc_text, Inches(0.94), top + Inches(0.12),
                W - Inches(1.3), desc_h - Inches(0.2), sz=14, col=self.P["text"])

            # DOMINANT visual — takes up 70% of slide height
            vt = top + desc_h + Inches(0.15)
        else:
            # No description text — give the full height to the visual
            vt = top
        vh = H - vt - Inches(0.3)
        all_paths = [d.get("image_path", "")] + d.get("extra_image_paths", []) if d.get("image_path") or d.get("extra_image_paths") else []
        self._premium_image_frame(slide, all_paths, Inches(0.25), vt,
                                   W - Inches(0.43), vh, d.get("visual_suggestion", "[ Architecture ]"),
                                   chart_data=d.get("chart_data"))


    # ── LAYOUT 5: TIMELINE ────────────────────────────────────────────────────
    def _lay_aesthetic_timeline(self, d: dict, cur: int, total: int):
        slide = self._blank()
        self._header(slide, d.get("title", ""))
        self._progress(slide, cur, total)

        nodes = d.get("nodes", [])
        if not nodes: return

        top   = Inches(0.9)
        avail = H - top - Inches(0.35)
        track_y = top + avail * 0.50

        # Track line
        _rect(slide, Inches(0.3), track_y - Inches(0.04),
              W - Inches(0.6), Inches(0.08), fill=self.P["border"])
        _rect(slide, Inches(0.3), track_y - Inches(0.02),
              W - Inches(0.6), Inches(0.04), fill=self.P["ac1"])

        nw = (W - Inches(0.6)) / max(len(nodes), 1)
        tag_cols = self.P["tag_colors"]

        for i, n in enumerate(nodes[:5]):
            ncx = Inches(0.3) + i * nw + nw / 2
            color = tag_cols[i % len(tag_cols)]

            # Node dot
            dot_r = Inches(0.22)
            _oval(slide, ncx - dot_r, track_y - dot_r, dot_r*2, dot_r*2,
                  fill=color, line=self.P["bg2"], lw=3.0)
            _oval(slide, ncx - dot_r*0.4, track_y - dot_r*0.4,
                  dot_r*0.8, dot_r*0.8, fill=self.P["bg"])

            # Card
            cw2 = nw - Inches(0.25)
            ch2 = avail * 0.42
            above = (i % 2 == 0)
            cy2 = track_y - ch2 - Inches(0.35) if above else track_y + Inches(0.38)

            # Connector line
            line_t = cy2 + ch2 if above else track_y + dot_r
            line_b = track_y - dot_r if above else cy2
            _rect(slide, ncx - Inches(0.03), min(line_t, line_b),
                  Inches(0.06), abs(line_b - line_t), fill=color)

            # Clean Single Card
            bw = getattr(self, "S", {}).get("card_border_width", 1.5)
            border_col = color if getattr(self, "S", {}).get("colored_borders", True) else self.P["border"]
            _round(slide, ncx - cw2/2, cy2, cw2, ch2, fill=self.P["card"], line=border_col, lw=bw)
            
            # Text formatting
            _tb(slide, n.get("header", "Phase").upper(), ncx - cw2/2 + Inches(0.12),
                cy2 + Inches(0.08), cw2 - Inches(0.24), Inches(0.35),
                sz=12, bold=True, col=color, align=PP_ALIGN.CENTER)
            
            # Use a slightly larger font and top align it to fill out the timeline card better
            _tb(slide, n.get("text", ""), ncx - cw2/2 + Inches(0.12),
                cy2 + Inches(0.45), cw2 - Inches(0.24), ch2 - Inches(0.55),
                sz=12.5, col=self.P["sub"], align=PP_ALIGN.LEFT)

    # ── LAYOUT 6: COMPARISON ──────────────────────────────────────────────────
    def _lay_aesthetic_comparison(self, d: dict, cur: int, total: int):
        slide = self._blank()
        self._header(slide, d.get("title", ""))
        self._progress(slide, cur, total)

        top = Inches(0.86)
        cd = d.get("chart_data")
        has_chart = False
        if cd and isinstance(cd, dict):
            # Check if it has the required data for a comparison chart, or if it's a fallback bar chart with values
            if (cd.get("left_values") and cd.get("right_values")) or cd.get("values"):
                has_chart = True
                
        # If chart_data available and valid, shrink columns to make room for chart at bottom
        col_area_h = (H - top - Inches(0.3)) * 0.55 if has_chart else H - top - Inches(0.3)
        avail = col_area_h
        cw2   = (W - Inches(0.7)) / 2

        for si, (hk, bk, color) in enumerate([
            ("left_header",  "left_bullets",  self.P["ac1"]),
            ("right_header", "right_bullets", self.P["ac2"]),
        ]):
            sx = Inches(0.25) + si * (cw2 + Inches(0.2))
            bw = getattr(self, "S", {}).get("card_border_width", 1.5)
            border_col = color if getattr(self, "S", {}).get("colored_borders", True) else self.P["border"]
            _round(slide, sx, top, cw2, avail, fill=self.P["card"], line=border_col, lw=bw)
            
            # Header line
            _tb(slide, d.get(hk, "Column").upper(), sx + Inches(0.15), top + Inches(0.11),
                cw2 - Inches(0.3), Inches(0.4), sz=14, bold=True,
                col=color, align=PP_ALIGN.CENTER)
            
            # Divider
            _rect(slide, sx + Inches(0.3), top + Inches(0.48), cw2 - Inches(0.6), Inches(0.02), fill=color)

            # Bullets
            blist = d.get(bk, [])
            if not blist:
                prefix = "left_" if "left" in bk else "right_"
                alt = d.get(f"{prefix}text", d.get(f"{prefix}description", ""))
                if alt:
                    blist = [alt]
                else:
                    blist = ["(No detailed content provided)"]
                    
            bsh = (avail - Inches(0.65)) / max(len(blist), 1)
            by = top + Inches(0.58)
            for bi, b in enumerate(blist[:5]):
                bold, txt = _parse_bullet(b)
                
                # Accent bar
                _rect(slide, sx + Inches(0.15), by + Inches(0.1), Inches(0.05), bsh - Inches(0.2), fill=color)
                # Bold sub-header
                if bold:
                    _tb(slide, bold.upper(), sx + Inches(0.3), by + Inches(0.08),
                        cw2 - Inches(0.45), Inches(0.3), sz=10.5, bold=True, col=color)
                    _tb(slide, txt, sx + Inches(0.3), by + Inches(0.3),
                        cw2 - Inches(0.45), bsh - Inches(0.36), sz=11, col=self.P["text"])
                else:
                    _tb(slide, txt, sx + Inches(0.3), by + Inches(0.08),
                        cw2 - Inches(0.45), bsh - Inches(0.16), sz=11.5, col=self.P["text"])
                by += bsh

        # VS badge centre
        vx = Inches(0.25) + cw2 + Inches(0.05)
        _round(slide, vx, top + avail/2 - Inches(0.5), Inches(0.1), Inches(1.0),
               fill=self.P["border"], line=None)
        _tb(slide, "VS", vx - Inches(0.06), top + avail/2 - Inches(0.22),
            Inches(0.22), Inches(0.45), sz=9, bold=True,
            col=self.P["text"], align=PP_ALIGN.CENTER)

        # ── Chart panel at bottom (if chart_data available) ────────────────
        if has_chart:
            chart_top = top + col_area_h + Inches(0.15)
            chart_h   = H - chart_top - Inches(0.3)
            all_paths = [d.get("image_path", "")] + d.get("extra_image_paths", []) if d.get("image_path") or d.get("extra_image_paths") else []
            self._premium_image_frame(
                slide, all_paths, Inches(0.25), chart_top,
                W - Inches(0.43), chart_h,
                d.get("visual_suggestion", "[ Comparison Chart ]"),
                chart_data=d.get("chart_data")
            )

    # ── LAYOUT 7: METRICS — KPI ROW + DOMINANT VISUAL ────────────────────────
    def _lay_aesthetic_metrics(self, d: dict, cur: int, total: int):
        slide = self._blank()
        self._header(slide, d.get("title", ""))
        self._progress(slide, cur, total)

        top     = Inches(0.86)
        metrics = d.get("metrics", [])
        n       = max(len(metrics[:4]), 1)
        gap     = Inches(0.18)
        mw      = (W - Inches(0.5) - gap*(n-1)) / n
        mh      = Inches(1.65)

        # KPI cards row — ROUNDED
        for i, m in enumerate(metrics[:4]):
            mx = Inches(0.25) + i*(mw + gap)
            color = self.P["tag_colors"][i % 4]
            bw = getattr(self, "S", {}).get("card_border_width", 1.5)
            border_col = color if getattr(self, "S", {}).get("colored_borders", True) else self.P["border"]
            _round(slide, mx, top, mw, mh, fill=self.P["card"], line=border_col, lw=bw)
            
            # Value
            _tb(slide, m.get("value", "—"), mx + Inches(0.1), top + Inches(0.15),
                mw - Inches(0.2), Inches(0.85), sz=34, bold=True,
                col=color, align=PP_ALIGN.CENTER)
            # Divider
            _rect(slide, mx + Inches(0.25), top + Inches(1.0), mw - Inches(0.5),
                  Inches(0.02), fill=color)
            # Label
            _tb(slide, m.get("label", ""), mx + Inches(0.12), top + Inches(1.2),
                mw - Inches(0.24), Inches(0.42), sz=9.5, col=self.P["sub"],
                align=PP_ALIGN.CENTER)

        # Dominant visual below (60%+ of slide) — auto-generate dashboard chart
        vt = top + mh + Inches(0.18)
        vh = H - vt - Inches(0.3)
        desc_fallback = d.get("description", d.get("text", ""))
        if not desc_fallback and "metrics" in d:
            desc_fallback = " ".join([m.get("label", "") if isinstance(m, dict) else str(m) for m in d["metrics"]])
        fallback_text = desc_fallback if desc_fallback else d.get("visual_suggestion", "[ Detailed Metric Analysis ]")

        # Auto-generate a bar chart from metrics data if no chart_data exists
        auto_chart = d.get("chart_data")
        if not auto_chart and metrics:
            auto_chart = {
                "type": "metrics", # Let ChartEngine handle extraction safely
                "title": d.get("title", "Key Metrics"),
                "metrics": metrics[:4]
            }

        self._premium_image_frame(slide, d.get("image_path", ""), Inches(0.25), vt,
                                   W - Inches(0.43), vh, fallback_text,
                                   chart_data=auto_chart)

    # ── LAYOUT 8: PITCH (DENSE GRID + VISUAL) ───────────────────────────────
    def _lay_aesthetic_pitch(self, d: dict, cur: int, total: int):
        slide = self._blank()
        self._header(slide, d.get("title", ""))
        self._progress(slide, cur, total)

        top = Inches(0.86)
        avail = H - top - Inches(0.3)
        
        # Top half: 4 dense grid cards
        top_h = Inches(2.3) # Increased from 1.3 to fit 40-60 words!
        cards = d.get("cards", [])
        if not cards and "bullets" in d:
            for i, b in enumerate(d["bullets"]):
                bold, text = _parse_bullet(b)
                header = bold if bold else f"Step {i+1}"
                cards.append({"header": header, "bullets": [text] if text else ["(No content)"]})
        
        # Parse any string-cards into dicts
        cards = [_parse_card(c) for c in cards]
        
        # Fix missing headers dynamically
        for i, c in enumerate(cards):
            if not c.get("header", ""):
                c["header"] = f"Module {i+1}"
                
        num_cards = max(len(cards[:4]), 1)
        total_gap = Inches(0.15) * (num_cards - 1)
        cw = (W - Inches(0.5) - total_gap) / num_cards
        
        for i, c in enumerate(cards[:4]):
            cx = Inches(0.25) + i*(cw + Inches(0.15))
            self._colored_card_full(slide, cx, top, cw, top_h, c, self.P["tag_colors"][i % 4])

        # Bottom half: Left image, Right bullets
        bot_t = top + top_h + Inches(0.2)
        bot_h = avail - top_h - Inches(0.2)
        
        bullets = d.get("right_bullets", d.get("bottom_bullets", []))
        if not bullets and "description" in d: bullets = [{"bold": "", "text": d["description"]}]
        
        # Check if we actually have meaningful bullets
        has_real_bullets = False
        for b in bullets:
            _, txt = _parse_bullet(b)
            if txt and "No detailed content" not in txt and txt.strip():
                has_real_bullets = True
                break

        # Left image (Dynamic width based on right content)
        lw = (W - Inches(0.7)) * 0.55 if has_real_bullets else (W - Inches(0.5))
        
        desc_fallback = d.get("description", d.get("text", ""))
        fallback_text = desc_fallback if desc_fallback else d.get("visual_suggestion", "[ Flowchart / Architecture ]")
        all_paths = [d.get("image_path", "")] + d.get("extra_image_paths", []) if d.get("image_path") or d.get("extra_image_paths") else []
        self._premium_image_frame(slide, all_paths, Inches(0.25), bot_t, lw, bot_h, fallback_text,
                                   chart_data=d.get("chart_data"))
        
        # If no meaningful bullets, skip drawing the right side cards
        if not has_real_bullets:
            return

        # Right bullets (45% width)
        rx = Inches(0.25) + lw + Inches(0.2)
        rw = (W - Inches(0.7)) * 0.45
        
        # Internal bullets
        by = bot_t
        # Cap bullet height so 1 bullet doesn't stretch 3.5 inches tall
        bh = min(bot_h / max(len(bullets), 1), Inches(0.85))
        
        for i, b in enumerate(bullets[:6]):
            bold, txt = _parse_bullet(b)
            if not bold:
                bold = f"Point {i+1}"
            self._bullet_card(slide, rx, by, rw, bh - Inches(0.1), bold, txt, self.P["tag_colors"][i%4], i+1)
            by += bh

    # ── LAYOUT: GALLERY — 1 to 6 images in an adaptive premium grid ───────────
    def _lay_aesthetic_gallery(self, d: dict, cur: int, total: int):
        """
        Pure visual layout for 1–6 images.
        
        Grid patterns:
          1 image  → full-width hero
          2 images → side-by-side equal panels
          3 images → large left + 2 stacked right
          4 images → 2×2 grid
          5 images → 2×2 top + 1 wide bottom
          6 images → 2×3 grid
        
        All images use contain-fit (never cropped or stretched).
        Titles and captions are derived from user hints stored in image_paths_hints.
        """
        slide = self._blank()
        self._header(slide, d.get("title", ""))
        self._progress(slide, cur, total)

        top     = Inches(0.86)
        avail_h = H - top - Inches(0.32)
        avail_w = W - Inches(0.5)
        lx      = Inches(0.25)
        gap     = Inches(0.15)

        # Collect all image paths for this gallery slide
        images: list[str] = []
        for p in d.get("image_paths", []):
            clean = _clean_image_path(str(p))
            if clean and Path(clean).exists():
                images.append(clean)
        # Also accept a single image_path
        if not images:
            single = _clean_image_path(d.get("image_path", ""))
            if single and Path(single).exists():
                images.append(single)

        n = len(images)

        def _frame(img, x, y, w, h):
            self._premium_image_frame(slide, img, x, y, w, h,
                                      d.get("visual_suggestion", ""))

        if n == 0:
            # Placeholder
            self._premium_image_frame(slide, "", lx, top, avail_w, avail_h,
                                      d.get("visual_suggestion", "[ Gallery ]"))

        elif n == 1:
            # Full-width hero
            _frame(images[0], lx, top, avail_w, avail_h)

        elif n == 2:
            # Side-by-side equal panels
            w = (avail_w - gap) / 2
            _frame(images[0], lx, top, w, avail_h)
            _frame(images[1], lx + w + gap, top, w, avail_h)

        elif n == 3:
            # Large left panel + 2 stacked on right
            main_w = avail_w * 0.58
            side_w = avail_w - main_w - gap
            side_h = (avail_h - gap) / 2
            _frame(images[0], lx,                 top,              main_w, avail_h)
            _frame(images[1], lx + main_w + gap,  top,              side_w, side_h)
            _frame(images[2], lx + main_w + gap,  top + side_h + gap, side_w, side_h)

        elif n == 4:
            # 2×2 perfect grid
            cw = (avail_w - gap) / 2
            ch = (avail_h - gap) / 2
            for idx, img in enumerate(images[:4]):
                col, row = idx % 2, idx // 2
                _frame(img, lx + col * (cw + gap), top + row * (ch + gap), cw, ch)

        elif n == 5:
            # Top 2×2 + bottom wide banner
            top_h  = avail_h * 0.54
            bot_h  = avail_h - top_h - gap
            cw     = (avail_w - gap) / 2
            for idx, img in enumerate(images[:4]):
                col, row = idx % 2, idx // 2
                _frame(img, lx + col * (cw + gap), top + row * (top_h / 2 + gap / 2), cw, top_h / 2 - gap / 2)
            _frame(images[4], lx, top + top_h + gap, avail_w, bot_h)

        else:  # 6+
            # 3-column × 2-row grid (shows first 6)
            cw = (avail_w - 2 * gap) / 3
            ch = (avail_h - gap) / 2
            for idx, img in enumerate(images[:6]):
                col, row = idx % 3, idx // 3
                _frame(img, lx + col * (cw + gap), top + row * (ch + gap), cw, ch)

def _auto_select_image_layout(slide_dict: dict) -> str:
    """
    After an image is assigned to a slide, pick the most aesthetically appropriate
    rendering layout based on actual content + image characteristics.

    Priority order:
    1. Gallery / showcase / poster are never overridden if explicitly set.
    2. Dense content (cards + bullets, or cards + image) → aesthetic_poster.
    3. Portrait images (tall) + text → aesthetic_flow.
    4. Rich bullets (3+) + landscape image → aesthetic_split.
    5. Cards only → aesthetic_grid.
    6. Sparse text + any image → aesthetic_flow.
    7. No text at all → aesthetic_showcase (let the image dominate).
    """
    layout = slide_dict.get("layout", "aesthetic_split")

    # Locked layouts — never override
    if layout in ("aesthetic_gallery", "aesthetic_showcase", "aesthetic_poster",
                  "aesthetic_title"):
        return layout

    bullets    = [b for b in slide_dict.get("bullets", []) if b]
    cards      = [c for c in slide_dict.get("cards",   []) if c]
    metrics    = slide_dict.get("metrics", [])
    img_aspect = float(slide_dict.get("image_aspect", 1.78))

    has_bullets   = len(bullets) >= 1
    has_rich_text = len(bullets) >= 3
    has_cards     = len(cards)   >= 1
    # Use image_slot flag (set by the pipeline) rather than checking file existence
    # because this function may run before image_path is applied to the slide dict
    has_img       = bool(slide_dict.get("image_slot") or slide_dict.get("image_path"))
    is_portrait   = img_aspect < 0.85

    # ── Dense content: route to poster ────────────────────────────────────────
    # The poster layout places cards at the top and text+image at the bottom.
    # We must only choose this if we actually have text (bullets or fallback description) 
    # to put next to the image. Otherwise, it leaves awkward empty space!
    has_poster_text = has_bullets or bool(slide_dict.get("description") or slide_dict.get("subtitle"))
    
    if has_cards and has_poster_text:
        return "aesthetic_poster"

    # ── Portrait image + text → flow ──────────────────────────────────────────
    if is_portrait and (has_bullets or has_cards):
        return "aesthetic_flow"

    # ── Bullets only + landscape image → split ────────────────────────────────
    if has_rich_text and not has_cards:
        return "aesthetic_split"

    # ── Metrics layout: keep as-is ────────────────────────────────────────────
    if metrics and layout == "aesthetic_metrics":
        return "aesthetic_metrics"

    # ── Cards only (no image) → standard grid ─────────────────────────────────
    if has_cards:
        return "aesthetic_grid"

    # ── Some bullets + landscape image ────────────────────────────────────────
    if has_bullets:
        return "aesthetic_split"

    # ── Nothing: full-screen showcase ─────────────────────────────────────────
    return "aesthetic_showcase"



def _normalize_and_recover(chunk_data, outline_chunk):
    """Maps hallucinated content arrays to expected keys and supplies fallbacks so no slide is ever blank.
    Image-only slides (no text) are NOT injected with 'Auto-Recovered' — the auto-layout
    selector will route them to aesthetic_showcase at build time.
    """
    if not isinstance(chunk_data, dict): return
    slides = chunk_data.get("slides", [])

    _VALID_LAYOUTS = frozenset([
        "aesthetic_split", "aesthetic_pitch", "aesthetic_flow", "aesthetic_grid",
        "aesthetic_timeline", "aesthetic_metrics", "aesthetic_comparison",
        "aesthetic_title", "aesthetic_showcase", "aesthetic_gallery",
        "aesthetic_poster",  # dynamic dashboard layout
    ])

    for i, s in enumerate(slides):
        expected_lay = outline_chunk[i].get("layout", "aesthetic_split") if i < len(outline_chunk) else "aesthetic_split"
        lay = s.get("layout", expected_lay)

        # 1. Force valid layout (handles hallucinated layout names)
        if lay not in _VALID_LAYOUTS:
            lay = expected_lay if expected_lay in _VALID_LAYOUTS else "aesthetic_split"
        s["layout"] = lay

        # 2. NEVER skip recovery based on image_slot.
        #    The AI often sets image_slot:true in chunk JSON but generates zero bullets/cards.
        #    We MUST always check and recover content regardless of image flags.
        #    (image_path is the ENGINE-set field; image_slot from AI JSON is just a hint.)

        # 3. Data Recovery for text-dependent layouts
        alt = (s.get("bullets") or s.get("cards") or s.get("nodes") or
               s.get("points") or s.get("items") or s.get("content") or s.get("metrics") or [])
        if isinstance(alt, str):
            alt = [{"bold": "Note", "text": alt}]

        if lay in ("aesthetic_split", "aesthetic_pitch", "aesthetic_flow"):
            if not s.get("bullets"):
                s["bullets"] = alt if alt else [{"bold": "Content", "text": "Content generation produced no detailed text for this section."}]
        elif lay in ("aesthetic_grid", "aesthetic_poster"):
            if not s.get("cards"):
                s["cards"] = alt if alt else [{"header": "Section", "bullets": ["Content missing"]}]
            # aesthetic_poster also needs bullets
            if lay == "aesthetic_poster" and not s.get("bullets"):
                s["bullets"] = [{"bold": "Key Point", "text": "Supporting detail for this section."}]
        elif lay == "aesthetic_timeline":
            if not s.get("nodes"):
                fallback_nodes = []
                for item in (alt or []):
                    bold, text = _parse_bullet(item) if not isinstance(item, dict) else (item.get("bold",""), item.get("text",""))
                    fallback_nodes.append({"header": bold or "Phase", "text": text or "Content missing"})
                s["nodes"] = fallback_nodes if fallback_nodes else [{"header": "Phase", "text": "Content missing"}]
        elif lay == "aesthetic_metrics":
            if not s.get("metrics"):
                s["metrics"] = alt if alt else [{"value": "-", "label": "Missing"}]
        elif lay == "aesthetic_comparison":
            if not s.get("left_bullets"):
                half = len(alt) // 2
                s["left_bullets"]  = alt[:half] if alt else [{"bold": "Left",  "text": "Content missing"}]
                s["right_bullets"] = alt[half:] if alt else [{"bold": "Right", "text": "Content missing"}]

def ppt_create(prompt: str, style: str = None, output_path: str = None,
               theme_image_path: str = None, research_data: dict = None,
               purpose: str = None, image_paths: list = None,
               image_descriptions: list = None):
    """
    End-to-end PPT generation pipeline.

    Args:
        image_descriptions: Parallel list of user-provided hints for each image in
                            image_paths (e.g. "product screenshot for the demo slide").
                            May contain empty strings. Used by ppt_image_engine to
                            semantically match each image to the best slide.
    """
    yield "🤖 Generating Presentation Outline...\n"
    
    if not purpose:
        purpose = _detect_purpose(prompt)

    # ── Extract explicit slide count from prompt ──────────────────────────────
    _slide_count_m = re.search(r'(\d+)\s*slides?', prompt, re.IGNORECASE)
    _target_slides = 0
    if _slide_count_m:
        _n = int(_slide_count_m.group(1))
        # Clamp between 4 and 20 for sanity
        _n = max(4, min(20, _n))
        _target_slides = _n
        slide_count_rule = f"Outline EXACTLY {_n} slides. No more, no less. STRICTLY {_n} slides."
    else:
        slide_count_rule = "Outline 8 to 12 slides."

    purpose_rules = ""
    if purpose == "hackathon":
        purpose_rules = "PURPOSE: This is a hackathon/startup pitch. Favor 'aesthetic_metrics', 'aesthetic_pitch', 'aesthetic_comparison' layouts (data-heavy, visual-dominant)."
    else:
        purpose_rules = "PURPOSE: This is a professional/general presentation. Favor 'aesthetic_split', 'aesthetic_grid', 'aesthetic_flow' layouts (text-balanced, clean)."

    # ── Validate user images and set outline rules ────────────────────────────
    valid_images = []
    if image_paths:
        for p in (image_paths or []):
            p = (p or "").strip()
            if p and Path(p).exists():
                valid_images.append(p)
    
    image_rules = ""
    if valid_images:
        num_imgs = len(valid_images)
        image_rules = (
            f"IMAGE INTEGRATION: The user has uploaded {num_imgs} image(s). "
            "Jarvis will automatically match each image to the most relevant slide based on "
            "the image descriptions the user provided — you do NOT need to assign images yourself. "
            "However, you MUST include enough slides with image-capable layouts so the images "
            "have beautiful homes. Use layouts like 'aesthetic_split' (text+image side-by-side), "
            "'aesthetic_showcase' (full-screen hero), or 'aesthetic_pitch' (impact grid + visual). "
            f"Include at least {min(num_imgs, 4)} slides with these image-capable layouts."
        )
    else:
        image_rules = "No user images were uploaded. Do not set 'image_slot' to true on any slide."


    # ── INJECT RESEARCH DATA (PHASE 3) ──
    enriched_prompt = prompt
    if research_data:
        facts_str = "\n- ".join(research_data.get("verified_facts", []))
        stats_str = "\n- ".join([f"{s['label']}: {s['value']}" for s in research_data.get("statistics", [])])
        sources_str = ", ".join(research_data.get("sources", []))
        
        enriched_prompt += f"""

==================================================
VERIFIED RESEARCH DATA (MUST USE THESE FACTS):
==================================================
FACTS:
- {facts_str}

STATISTICS & METRICS:
- {stats_str}

SOURCES (Cite these if applicable): {sources_str}
==================================================
"""
        yield f"📚 Injected {len(research_data.get('verified_facts', []))} verified facts from live web research.\n"

    raw_outline = _groq_call(_SYS_OUTLINE, _USR_OUTLINE.format(
        prompt=enriched_prompt,
        purpose_rules=purpose_rules,
        slide_count_rule=slide_count_rule,
        image_rules=image_rules,
    ))
    plan = _parse(raw_outline)
    plan["purpose"] = purpose
    
    if style and style in PERSONALITIES:
        plan["personality"] = style

    # If user uploaded a theme reference image, extract its colors and use them
    if theme_image_path:
        yield "🎨 Extracting color theme from your reference image...\n"
        try:
            extracted = extract_theme_from_image(theme_image_path)
            plan["custom_theme"] = extracted
            yield f"✅ Theme extracted! Primary accent: #{extracted.get('ac1','?')}\n"
        except Exception as e:
            yield f"⚠️ Could not extract theme from image ({e}). Using auto-theme.\n"
    slides_outline = plan.get("slides", [])
    
    # ── ALGORITHMIC SLIDE LIMIT ENFORCEMENT ──
    # If the user explicitly requested N slides, and the LLM ignored the prompt 
    # rule by generating > N slides, we strictly truncate the outline.
    if _target_slides > 0 and len(slides_outline) > _target_slides:
        yield f"⚠️ Enforcing strict slide limit: Truncating from {len(slides_outline)} to {_target_slides} slides.\n"
        # Preserve the very last slide (it's usually the Conclusion/Thank You)
        last_slide = slides_outline[-1]
        slides_outline = slides_outline[:_target_slides - 1]
        last_slide["slide_number"] = _target_slides
        slides_outline.append(last_slide)
        
    total_slides = len(slides_outline)
    yield f"\U0001f4cb Outline created: {total_slides} slides planned.\n"

    if valid_images:
        yield f"\U0001f5bc\ufe0f {len(valid_images)} user image(s) detected \u2014 mapping to layout slots.\n"
    
    full_slides = []
    
    # Chunk generation (2 slides at a time to prevent hallucination)
    chunk_size = 2
    outline_str = json.dumps([{"slide_number": s.get("slide_number"), "title": s.get("title"), "layout": s.get("layout")} for s in slides_outline])

    for i in range(0, total_slides, chunk_size):
        chunk = slides_outline[i:i+chunk_size]
        start_idx = chunk[0].get('slide_number', i+1)
        end_idx = chunk[-1].get('slide_number', i+len(chunk))
        
        yield f"✍️ Generating extremely dense content for Slides {start_idx}-{end_idx}...\n"
        
        chunk_slides = []
        last_err = None
        for attempt in range(3):  # Up to 3 attempts with exponential backoff
            try:
                raw_chunk = _groq_call(
                    _SYS_CHUNK,
                    _USR_CHUNK.format(prompt=enriched_prompt, outline=outline_str, start_idx=start_idx, end_idx=end_idx)
                )
                chunk_data = _parse(raw_chunk)
                _normalize_and_recover(chunk_data, chunk)
                chunk_slides = chunk_data.get("slides", [])
                # After _normalize_and_recover has already fixed empty slides,
                # just check we got at least one slide so we can map it
                if len(chunk_slides) > 0:
                    break  # Success
                last_err = ValueError(f"Got 0 slides, expected {len(chunk)}.")
            except Exception as e:
                last_err = e
            if attempt < 2:  # Don't sleep after last attempt
                import time as _t; _t.sleep(1.5 * (2 ** attempt))

        if chunk_slides:
            # Enforce strict mapping to prevent LLM from hallucinating extra slides
            for i, expected_c in enumerate(chunk):
                expected_num = expected_c.get('slide_number')
                # 1. Try finding by exact slide number
                match_idx = next((idx for idx, s in enumerate(chunk_slides) if s.get('slide_number') == expected_num), None)
                # 2. Fallback to positional index if not found by number (always pick the first remaining one)
                if match_idx is None and len(chunk_slides) > 0:
                    match_idx = 0
                
                if match_idx is not None:
                    # Pop the slide so it cannot be matched again for a later slide in the chunk
                    match = chunk_slides.pop(match_idx)
                    # Guarantee slide number and fallback layout matches outline expectation
                    match['slide_number'] = expected_num
                    if 'layout' not in match:
                        match['layout'] = expected_c.get('layout', 'aesthetic_split')
                    full_slides.append(match)
                else:
                    # Fallback for missing slide
                    full_slides.append({
                        "slide_number": expected_num,
                        "title": expected_c.get("title", "Slide"),
                        "layout": expected_c.get("layout", "aesthetic_split"),
                        "bullets": [{"bold": expected_c.get("title", "Note"), "text": "Content could not be generated for this slide."}]
                    })

        else:
            yield f"⚠️ Chunk {start_idx}-{end_idx} failed after 2 attempts ({last_err}). Using fallback.\n"
            for c in chunk:
                full_slides.append({
                    "slide_number": c.get("slide_number"),
                    "title": c.get("title", "Slide"),
                    "layout": c.get("layout", "aesthetic_split"),
                    "bullets": [{"bold": c.get("title", "Note"), "text": "Content could not be generated for this slide."}]
                })
    yield "🎨 Content generated! Extracting chart data for infographics...\n"

    # ── NEW: Chart Data Extraction Pass ──────────────────────────────────
    # For each slide with a visual_suggestion, ask the LLM to extract
    # structured chart data so ChartEngine can render real infographics.
    chart_ok = 0
    chart_fail = 0
    used_chart_types = []
    for si, sd in enumerate(full_slides):
        vs = sd.get("visual_suggestion", "")
        lay = sd.get("layout", "")
        # Skip title slides — they don't need charts
        if lay == "aesthetic_title" or not vs:
            continue
        # Build a content summary for context
        content_parts = []
        for b in sd.get("bullets", []):
            if isinstance(b, dict):
                content_parts.append(f"{b.get('bold','')} {b.get('text','')}".strip())
        for c in sd.get("cards", []):
            if isinstance(c, dict):
                content_parts.append(c.get("header", ""))
        for m in sd.get("metrics", []):
            if isinstance(m, dict):
                content_parts.append(f"{m.get('value','')} {m.get('label','')}".strip())
        for n in sd.get("nodes", []):
            if isinstance(n, dict):
                content_parts.append(f"{n.get('header','')} {n.get('text','')}".strip())
        content_summary = "; ".join(content_parts[:6]) or sd.get("title", "")

        try:
            raw_chart = _groq_call(
                _SYS_CHART,
                _USR_CHART.format(
                    title=sd.get("title", ""),
                    layout=lay,
                    visual_suggestion=vs,
                    content_summary=content_summary[:500],
                    used_charts=", ".join(used_chart_types) if used_chart_types else "None"
                ),
                tokens=800
            )
            chart_data = _parse(raw_chart)
            chart_type = chart_data.get("type", "").lower() if chart_data else ""
            if chart_data and chart_type and chart_type != "none":
                sd["chart_data"] = chart_data
                used_chart_types.append(chart_data.get("type"))
                chart_ok += 1
        except Exception as e:
            chart_fail += 1
            # Non-fatal — slide will just show placeholder text
            print(f"[ppt] chart extraction failed for slide {si+1}: {e}")

    if chart_ok:
        yield f"📊 Extracted chart data for {chart_ok} slides.\n"
    if chart_fail:
        yield f"⚠️ {chart_fail} slides will use text placeholders (chart extraction failed).\n"

    yield "🎨 Building PPTX with precision curves + infographics...\n"

    # ── CRITICAL: Write the fully-populated slides back into the plan ──
    plan["slides"] = full_slides

    # ── Smart Image Distribution (via ppt_image_engine) ──────────────
    if valid_images:
        try:
            from app.services.ppt_image_engine import build_image_descriptors, match_images_to_slides

            yield f"\U0001f9e0 Analysing {len(valid_images)} image(s) — matching to slides...\n"

            # Build descriptors (user hints + filenames — zero extra API calls)
            descriptors = build_image_descriptors(valid_images, image_descriptions or [])

            # Semantically match images to slides
            image_map, overflow_images = match_images_to_slides(descriptors, full_slides)

            # ── Apply matches + auto-select best layout per slide ────────────────
            assigned = 0
            for slide_idx, desc in image_map.items():
                sd = full_slides[slide_idx]
                sd["image_path"]   = desc.path
                sd["image_aspect"]  = desc.aspect_ratio
                sd["image_slot"]    = True
                # Auto-select the best rendering layout based on content + image
                sd["layout"] = _auto_select_image_layout(sd)
                assigned += 1
                # Count extra images that were co-assigned to the same slide
                extra = sd.get("extra_image_paths", [])
                if extra:
                    assigned += len(extra)

            yield f"\U0001f5bc\ufe0f Matched {assigned} image(s) to best-fit slides with adaptive layouts.\n"

            # ── Overflow: group into gallery slides (up to 4 images each) ─────────
            if overflow_images:
                yield f"\U0001f5bc\ufe0f {len(overflow_images)} overflow image(s) → creating gallery slide(s)...\n"
                insert_pos = max(len(full_slides) - 1, 1)
                # Group every 4 overflow images into a single gallery slide
                batch_size = 4
                for batch_start in range(0, len(overflow_images), batch_size):
                    batch = overflow_images[batch_start : batch_start + batch_size]
                    # Build a title from the hints / filenames in this batch
                    titles = [
                        (desc.hint.strip() or Path(desc.path).stem.replace("_", " ").replace("-", " ").title())
                        for desc in batch
                    ]
                    gallery_title = " · ".join(t[:20] for t in titles[:3])
                    if len(titles) > 3:
                        gallery_title += f" +{len(titles)-3} more"

                    new_slide = {
                        "slide_number":   insert_pos + 1,
                        "layout":         "aesthetic_gallery",
                        "title":          gallery_title[:60] or "Visual Gallery",
                        "image_paths":    [desc.path for desc in batch],
                        "image_slot":     True,
                    }
                    full_slides.insert(insert_pos, new_slide)
                    insert_pos += 1
                    assigned += len(batch)

            # Re-number slides after insertions
            for idx, slide in enumerate(full_slides):
                slide["slide_number"] = idx + 1

            yield f"\U0001f5bc\ufe0f Successfully integrated {assigned} image(s) into the presentation.\n"

        except Exception as img_err:
            yield f"\u26a0\ufe0f Image matching failed ({img_err}). Using sequential fallback.\n"
            # Graceful fallback: sequential without layout-overwriting
            assigned = 0
            eligible = [i for i in range(1, len(full_slides) - 1)
                        if not full_slides[i].get("image_path")]
            leftover_for_gallery: list = []
            for img_path in valid_images:
                if eligible:
                    idx = eligible.pop(0)
                    aspect = _get_image_aspect_ratio(img_path)
                    full_slides[idx]["image_path"]   = img_path
                    full_slides[idx]["image_slot"]    = True
                    full_slides[idx]["image_aspect"]  = aspect
                    full_slides[idx]["layout"] = _auto_select_image_layout(full_slides[idx])
                    assigned += 1
                else:
                    leftover_for_gallery.append(img_path)

            # Group leftover images into gallery slides
            insert_pos = max(len(full_slides) - 1, 1)
            for batch_start in range(0, len(leftover_for_gallery), 4):
                batch_paths = leftover_for_gallery[batch_start : batch_start + 4]
                full_slides.insert(insert_pos, {
                    "slide_number": insert_pos + 1,
                    "layout":       "aesthetic_gallery",
                    "title":        "Visual Gallery",
                    "image_paths":  batch_paths,
                    "image_slot":   True,
                })
                insert_pos += 1
                assigned += len(batch_paths)

            for idx, slide in enumerate(full_slides):
                slide["slide_number"] = idx + 1
            yield f"\U0001f5bc\ufe0f Integrated {assigned} image(s) (fallback mode).\n"
    
    out = output_path or str(Path.home()/"Desktop"/f"Aesthetic_Deck_{int(time.time())}.pptx")
    b = PresentationBuilder(plan, out)
    for s in b.build_with_progress(): yield s + "\n"
    yield f"✅ Saved to: `{out}`\n"

def ppt_styles():
    return {"styles": {k: {"name": v["name"], "desc": v["desc"]} for k, v in PERSONALITIES.items()},
            "count": len(PERSONALITIES)}

def _pick(p): return "ocean_pro"
